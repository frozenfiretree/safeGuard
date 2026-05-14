import hashlib
import csv
import json
import re
import shutil
import subprocess
import tempfile
import time
import unicodedata
import uuid
import zipfile
from difflib import SequenceMatcher
from xml.etree import ElementTree as ET
from pathlib import Path
from typing import Any, Optional

from config_app import DATA_DIR
from detection.parsers import extract_file_content
from models import (
    FileRecord,
    FileVersion,
    ParseResult,
    TrackedFile,
    TrackedFileEvent,
    TrackedFileVersion,
)
from path_utils import remote_path_name
from storage import db_session, object_storage


TRACK_ROOT = DATA_DIR / "guard_state"
CONTENT_EVENTS = {"initial", "modified", "restored"}
RETAINED_CONTENT_SNAPSHOTS = 3


def _now() -> float:
    return time.time()


def _safe_name(value: str) -> str:
    return re.sub(r"[^A-Za-z0-9_.-]+", "_", value or "file")[:160]


def _norm_path(value: Optional[str]) -> str:
    return str(value or "").replace("/", "\\").strip().lower()


def _is_download_temp_path(value: Optional[str]) -> bool:
    name = Path(str(value or "")).name.lower()
    return (
        name.endswith((".crdownload", ".tmp", ".temp"))
        or name.startswith("~$")
        or name.startswith("~lock.")
        or bool(re.match(r"^未确认\s+\d+\.crdownload$", name, flags=re.IGNORECASE))
    )


def _hash_text(value: str) -> str:
    return hashlib.sha256(value.encode("utf-8", errors="ignore")).hexdigest()


def _file_key(agent_id: str, path: str, content_hash: str = "") -> str:
    seed = f"{agent_id}|{_norm_path(path)}"
    return _hash_text(seed)[:32]


def _tracked_dir(agent_id: str, tracked_file_id: str) -> Path:
    return TRACK_ROOT / agent_id / "files" / tracked_file_id


def _ensure_dirs(agent_id: str, tracked_file_id: str) -> dict[str, Path]:
    root = _tracked_dir(agent_id, tracked_file_id)
    dirs = {
        "root": root,
        "events": root / "events",
        "versions": root / "versions",
        "highlights": root / "highlights",
        "diffs": root / "diffs",
    }
    for path in dirs.values():
        path.mkdir(parents=True, exist_ok=True)
    return dirs


def _write_json(path: Path, payload: dict) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")


def _read_text_blocks(path: Optional[Path]) -> list[str]:
    if not path or not path.exists():
        return []
    result = extract_file_content(path)
    blocks = []
    for item in result.get("text_blocks") or []:
        text = str(item.get("text") or "").strip()
        if text:
            blocks.append(text)
    return blocks


def _diff_texts(old_path: Optional[Path], new_path: Optional[Path], rename: Optional[dict] = None, deleted: Optional[dict] = None) -> dict:
    old_texts = _read_text_blocks(old_path)
    new_texts = _read_text_blocks(new_path)
    matcher = SequenceMatcher(a=old_texts, b=new_texts, autojunk=False)
    added: list[str] = []
    removed: list[str] = []
    modified: list[dict[str, str]] = []
    for tag, i1, i2, j1, j2 in matcher.get_opcodes():
        if tag == "equal":
            continue
        if tag == "insert":
            added.extend(new_texts[j1:j2])
        elif tag == "delete":
            removed.extend(old_texts[i1:i2])
        elif tag == "replace":
            before = "\n".join(old_texts[i1:i2]).strip()
            after = "\n".join(new_texts[j1:j2]).strip()
            modified.append({"before": before, "after": after})
    summary_parts = []
    if added:
        summary_parts.append(f"新增 {len(added)} 处文本")
    if removed:
        summary_parts.append(f"删除 {len(removed)} 处文本")
    if modified:
        summary_parts.append(f"修改 {len(modified)} 个文本块")
    if rename:
        summary_parts.append(f"文件已重命名：{rename.get('old_name')} -> {rename.get('new_name')}")
    if deleted:
        summary_parts.append(f"文件已于 {time.strftime('%Y-%m-%d %H:%M:%S', time.localtime(float(deleted.get('deleted_at') or _now())))} 删除")
    if not summary_parts:
        summary_parts.append("文件内容无可解析文本差异")
    return {
        "summary": "，".join(summary_parts),
        "added_texts": added[:50],
        "removed_texts": removed[:50],
        "modified_blocks": modified[:50],
        "rename": rename,
        "deleted": deleted,
    }


def _extract_hits(parse_data: dict | None) -> list[dict[str, Any]]:
    hits = []
    for key in ("rule_findings", "ocr_findings", "llm_findings"):
        for item in (parse_data or {}).get(key) or []:
            matched = str(item.get("matched_text") or item.get("text") or "").strip()
            if not matched:
                continue
            hit = dict(item)
            hit["matched_text"] = matched
            hit["source_group"] = key
            hits.append(hit)
    return hits


def _hit_source_label(source_group: str) -> str:
    return {
        "rule_findings": "规则",
        "ocr_findings": "OCR",
        "llm_findings": "LLM",
    }.get(source_group, source_group or "未知")


def _summarize_detection_hits(hits: list[dict[str, Any]], parse_data: dict | None = None) -> dict:
    groups: dict[str, dict[str, Any]] = {}
    for hit in hits:
        group = str(hit.get("source_group") or "unknown")
        bucket = groups.setdefault(group, {"count": 0, "examples": []})
        bucket["count"] += 1
        text = str(hit.get("matched_text") or "").strip()
        location = str(hit.get("location") or "").strip()
        if text and len(bucket["examples"]) < 5:
            bucket["examples"].append({"text": text[:120], "location": location})
    parts = []
    for group in ("rule_findings", "ocr_findings", "llm_findings"):
        item = groups.get(group)
        if item and item["count"]:
            parts.append(f"{_hit_source_label(group)}命中 {item['count']} 处")
    llm_summary = str((parse_data or {}).get("llm_summary") or "").strip()
    if llm_summary and not groups.get("llm_findings"):
        parts.append(f"LLM摘要：{llm_summary[:160]}")
    return {
        "summary": "；".join(parts),
        "groups": groups,
        "llm_summary": llm_summary,
        "llm_used": bool((parse_data or {}).get("llm_used")),
        "llm_gate_reason": (parse_data or {}).get("llm_gate_reason"),
    }


def _merge_summary_parts(*parts: Optional[str]) -> str:
    seen = set()
    values = []
    for part in parts:
        text = str(part or "").strip()
        if not text or text in seen:
            continue
        seen.add(text)
        values.append(text)
    return "；".join(values)


def _hit_texts(hits: list[dict[str, Any]]) -> list[str]:
    seen = set()
    values = []
    for hit in hits:
        text = str(hit.get("matched_text") or "").strip()
        if not text or text in seen:
            continue
        seen.add(text)
        values.append(text)
    return values


WORD_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
ET.register_namespace("w", WORD_NS)


def _word_tag(name: str) -> str:
    return f"{{{WORD_NS}}}{name}"


def _highlight_docx_xml(source: Path, target: Path, hits: list[str]) -> Optional[str]:
    if not hits:
        return None
    xml_parts = {
        "word/document.xml",
        "word/footnotes.xml",
        "word/endnotes.xml",
    }
    try:
        with zipfile.ZipFile(source, "r") as zin:
            xml_parts.update(
                name
                for name in zin.namelist()
                if re.match(r"word/(header|footer)\d+\.xml$", name)
            )
            changed_parts: dict[str, bytes] = {}
            applied = 0
            for name in xml_parts:
                if name not in zin.namelist():
                    continue
                root = ET.fromstring(zin.read(name))
                part_applied = 0
                for paragraph in root.iter(_word_tag("p")):
                    paragraph_text = "".join(text.text or "" for text in paragraph.iter(_word_tag("t")))
                    if not _hit_matches(paragraph_text, hits):
                        continue
                    for run in paragraph.iter(_word_tag("r")):
                        run_text = "".join(text.text or "" for text in run.iter(_word_tag("t")))
                        if not run_text:
                            continue
                        if not _hit_matches(run_text, hits) and not any(hit in paragraph_text for hit in hits):
                            continue
                        rpr = run.find(_word_tag("rPr"))
                        if rpr is None:
                            rpr = ET.Element(_word_tag("rPr"))
                            run.insert(0, rpr)
                        highlight = rpr.find(_word_tag("highlight"))
                        if highlight is None:
                            highlight = ET.SubElement(rpr, _word_tag("highlight"))
                        highlight.set(_word_tag("val"), "yellow")
                        part_applied += 1
                if part_applied:
                    changed_parts[name] = ET.tostring(root, encoding="utf-8", xml_declaration=True)
                    applied += part_applied
            if applied <= 0:
                return None
            target.parent.mkdir(parents=True, exist_ok=True)
            with zipfile.ZipFile(target, "w", compression=zipfile.ZIP_DEFLATED) as zout:
                for item in zin.infolist():
                    data = changed_parts.get(item.filename)
                    if data is None:
                        data = zin.read(item.filename)
                    zout.writestr(item, data)
            return "docx_xml_highlight"
    except Exception:
        return None


def _highlight_docx(source: Path, target: Path, hits: list[str]) -> Optional[str]:
    if not hits:
        return None
    artifact_type = _highlight_docx_xml(source, target, hits)
    if artifact_type:
        return artifact_type
    from docx import Document
    from docx.enum.text import WD_COLOR_INDEX

    doc = Document(source)

    def mark_paragraph(paragraph):
        text = paragraph.text or ""
        if not any(hit in text for hit in hits):
            return
        for run in paragraph.runs:
            run_text = run.text or ""
            if any(hit in run_text for hit in hits) or run_text:
                run.font.highlight_color = WD_COLOR_INDEX.YELLOW

    for paragraph in doc.paragraphs:
        mark_paragraph(paragraph)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    mark_paragraph(paragraph)
    target.parent.mkdir(parents=True, exist_ok=True)
    doc.save(target)
    return "docx_highlight"


def _highlight_doc(source: Path, target: Path, hits: list[str]) -> tuple[Optional[str], str]:
    if not hits:
        return None, "no_hits"
    temp_dir = Path(tempfile.mkdtemp(prefix="safeguard_doc_highlight_"))
    try:
        converted = _convert_office_file(source, temp_dir, "docx")
        if not converted:
            return None, "doc_to_docx_conversion_failed"
        temp_highlight = temp_dir / f"{source.stem}_highlight.docx"
        artifact_type = _highlight_docx(converted, temp_highlight, hits)
        if not artifact_type:
            return None, "no_hits"
        converted_back = _convert_office_file(temp_highlight, temp_dir, "doc")
        if converted_back and converted_back.exists():
            target.parent.mkdir(parents=True, exist_ok=True)
            shutil.copyfile(converted_back, target)
            return "doc_highlight", "doc_to_docx_highlight_to_doc"
        fallback = target.with_suffix(".docx")
        target.parent.mkdir(parents=True, exist_ok=True)
        shutil.copyfile(temp_highlight, fallback)
        return "doc_converted_docx_highlight", "doc_to_docx_highlight"
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)


def _normalize_pdf_text(value: str) -> str:
    value = unicodedata.normalize("NFKC", value or "")
    return re.sub(r"\s+", "", value)


def _highlight_pdf(source: Path, target: Path, hits: list[str]) -> tuple[Optional[str], int]:
    if not hits:
        return None, 0
    import fitz

    doc = fitz.open(source)
    applied = 0
    normalized_hits = [(hit, _normalize_pdf_text(hit)) for hit in hits if hit.strip()]
    for page in doc:
        page_text = page.get_text("text") or ""
        searchable = bool(page_text.strip())
        for raw_hit, normalized_hit in normalized_hits:
            rects = page.search_for(raw_hit)
            if not rects and searchable and normalized_hit and normalized_hit in _normalize_pdf_text(page_text):
                # PyMuPDF search_for keeps native PDF text coordinates. If the exact
                # spacing differs, report fallback instead of pretending precision.
                continue
            for rect in rects:
                annot = page.add_highlight_annot(rect)
                if annot:
                    annot.set_info(content=f"Sensitive hit: {raw_hit[:80]}")
                    annot.update()
                    applied += 1
    if applied <= 0:
        doc.close()
        return None, 0
    target.parent.mkdir(parents=True, exist_ok=True)
    doc.save(target, garbage=4, deflate=True)
    doc.close()
    return "pdf_native_highlight_pdf", applied


def _classify_pdf_file(source: Path) -> dict:
    try:
        import fitz

        doc = fitz.open(source)
        pages_with_text = 0
        pages_without_text = 0
        image_blocks = 0
        for page in doc:
            if (page.get_text("text") or "").strip():
                pages_with_text += 1
            else:
                pages_without_text += 1
            image_blocks += len(page.get_images(full=True))
        page_count = len(doc)
        doc.close()
        if pages_with_text and pages_without_text:
            pdf_type = "hybrid_pdf"
        elif pages_with_text:
            pdf_type = "text_pdf"
        elif image_blocks:
            pdf_type = "scanned_pdf"
        else:
            pdf_type = "empty_pdf"
        return {
            "pdf_type": pdf_type,
            "page_count": page_count,
            "pages_with_text": pages_with_text,
            "pages_without_text": pages_without_text,
            "image_blocks": image_blocks,
        }
    except Exception as exc:
        return {"pdf_type": "unknown_pdf", "error": str(exc)}


def _parse_page_from_location(location: str) -> Optional[int]:
    match = re.search(r"(?:^|:)page:(\d+)(?:$|:)", location or "")
    if not match:
        return None
    return max(int(match.group(1)) - 1, 0)


def _highlight_pdf_ocr_boxes(source: Path, target: Path, hits: list[dict[str, Any]]) -> tuple[Optional[str], int]:
    import fitz

    doc = fitz.open(source)
    applied = 0
    for hit in hits:
        bbox = hit.get("bbox")
        if not bbox:
            continue
        page_idx = _parse_page_from_location(str(hit.get("location") or ""))
        if page_idx is None or page_idx >= len(doc):
            continue
        points = []
        for item in bbox:
            if isinstance(item, (list, tuple)) and len(item) >= 2:
                try:
                    points.append((float(item[0]), float(item[1])))
                except Exception:
                    continue
        if len(points) < 2:
            continue
        page = doc[page_idx]
        page_rect = page.rect
        xs = [point[0] for point in points]
        ys = [point[1] for point in points]
        max_x = max(xs) or page_rect.width
        max_y = max(ys) or page_rect.height
        scale_x = page_rect.width / max_x if max_x > page_rect.width * 1.25 else 1.0
        scale_y = page_rect.height / max_y if max_y > page_rect.height * 1.25 else 1.0
        rect = fitz.Rect(min(xs) * scale_x, min(ys) * scale_y, max(xs) * scale_x, max(ys) * scale_y)
        rect = rect & page_rect
        if rect.is_empty or rect.width < 1 or rect.height < 1:
            continue
        annot = page.add_rect_annot(rect)
        if annot:
            annot.set_colors(stroke=(1, 0.2, 0.1), fill=(1, 0.95, 0.2))
            annot.set_opacity(0.28)
            annot.set_info(content=f"Sensitive OCR hit: {str(hit.get('matched_text') or '')[:80]}")
            annot.update()
            applied += 1
    if applied <= 0:
        doc.close()
        return None, 0
    target.parent.mkdir(parents=True, exist_ok=True)
    doc.save(target, garbage=4, deflate=True)
    doc.close()
    return "pdf_ocr_bbox_highlight_pdf", applied


def _hit_matches(text: str, hits: list[str]) -> bool:
    value = str(text or "")
    if not value:
        return False
    normalized = _normalize_pdf_text(value)
    return any(hit in value or _normalize_pdf_text(hit) in normalized for hit in hits)


def _convert_office_file(source: Path, target_dir: Path, target_extension: str) -> Optional[Path]:
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        return None
    target_extension = target_extension.lower().lstrip(".")
    target_dir.mkdir(parents=True, exist_ok=True)
    try:
        proc = subprocess.run(
            [
                soffice,
                "--headless",
                "--convert-to",
                target_extension,
                "--outdir",
                str(target_dir),
                str(source),
            ],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            errors="replace",
            timeout=120,
        )
    except Exception:
        return None
    if proc.returncode != 0:
        return None
    expected = target_dir / f"{source.stem}.{target_extension}"
    if expected.exists():
        return expected
    matches = list(target_dir.glob(f"{source.stem}.*"))
    return matches[0] if matches else None


def _highlight_pptx(source: Path, target: Path, hits: list[str]) -> tuple[Optional[str], int]:
    if not hits:
        return None, 0
    from pptx import Presentation
    from pptx.dml.color import RGBColor

    prs = Presentation(source)
    applied = 0

    def mark_shape(shape) -> None:
        nonlocal applied
        if getattr(shape, "has_text_frame", False) and _hit_matches(getattr(shape, "text", ""), hits):
            try:
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(255, 242, 128)
            except Exception:
                pass
            try:
                shape.line.color.rgb = RGBColor(255, 64, 64)
            except Exception:
                pass
            applied += 1
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    if _hit_matches(cell.text, hits):
                        try:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(255, 242, 128)
                        except Exception:
                            pass
                        applied += 1
        if hasattr(shape, "shapes"):
            for child in shape.shapes:
                mark_shape(child)

    for slide in prs.slides:
        for shape in slide.shapes:
            mark_shape(shape)
    if applied <= 0:
        return None, 0
    target.parent.mkdir(parents=True, exist_ok=True)
    prs.save(target)
    return "pptx_shape_highlight", applied


def _highlight_ppt(source: Path, target: Path, hits: list[str]) -> tuple[Optional[str], int, str]:
    if not hits:
        return None, 0, "no_hits"
    temp_dir = Path(tempfile.mkdtemp(prefix="safeguard_ppt_highlight_"))
    try:
        converted = _convert_office_file(source, temp_dir, "pptx")
        if not converted:
            return None, 0, "ppt_to_pptx_conversion_failed"
        temp_highlight = temp_dir / f"{source.stem}_highlight.pptx"
        artifact_type, applied = _highlight_pptx(converted, temp_highlight, hits)
        if not artifact_type:
            return None, applied, "no_hits"
        converted_back = _convert_office_file(temp_highlight, temp_dir, "ppt")
        if converted_back and converted_back.exists():
            target.parent.mkdir(parents=True, exist_ok=True)
            shutil.copyfile(converted_back, target)
            return "ppt_shape_highlight", applied, "ppt_to_pptx_highlight_to_ppt"
        fallback = target.with_suffix(".pptx")
        target.parent.mkdir(parents=True, exist_ok=True)
        shutil.copyfile(temp_highlight, fallback)
        return "ppt_converted_pptx_shape_highlight", applied, "ppt_to_pptx_highlight"
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)


def _highlight_xlsx(source: Path, target: Path, hits: list[str]) -> tuple[Optional[str], int]:
    if not hits:
        return None, 0
    import openpyxl
    from openpyxl.styles import PatternFill

    wb = openpyxl.load_workbook(source)
    fill = PatternFill(fill_type="solid", fgColor="FFF280")
    applied = 0
    for sheet in wb.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if cell.value is not None and _hit_matches(str(cell.value), hits):
                    cell.fill = fill
                    applied += 1
    if applied <= 0:
        return None, 0
    target.parent.mkdir(parents=True, exist_ok=True)
    wb.save(target)
    return "xlsx_cell_highlight", applied


def _highlight_csv(source: Path, target: Path, hits: list[str]) -> tuple[Optional[str], int]:
    if not hits:
        return None, 0
    import openpyxl
    from openpyxl.styles import PatternFill

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "csv_highlight"
    fill = PatternFill(fill_type="solid", fgColor="FFF280")
    applied = 0
    with open(source, "r", encoding="utf-8-sig", errors="ignore", newline="") as handle:
        reader = csv.reader(handle)
        for row_idx, row in enumerate(reader, start=1):
            for col_idx, value in enumerate(row, start=1):
                cell = ws.cell(row=row_idx, column=col_idx, value=value)
                if _hit_matches(value, hits):
                    cell.fill = fill
                    applied += 1
    if applied <= 0:
        return None, 0
    target.parent.mkdir(parents=True, exist_ok=True)
    wb.save(target)
    return "csv_xlsx_highlight", applied


def _create_highlight(source: Optional[Path], target_dir: Path, version_no: int, name: str, hits: list[dict[str, Any]]) -> tuple[Optional[str], Optional[str], dict]:
    if not source or not source.exists():
        return None, None, {"status": "missing_source"}
    suffix = source.suffix.lower()
    hit_values = _hit_texts(hits)
    if suffix == ".docx":
        target = target_dir / f"v{version_no}_{_safe_name(Path(name).stem)}_highlight.docx"
        artifact_type = _highlight_docx(source, target, hit_values)
        return (str(target) if artifact_type else None, artifact_type, {"status": "ok" if artifact_type else "no_hits", "strategy": "docx_run_or_paragraph"})
    if suffix == ".doc":
        target = target_dir / f"v{version_no}_{_safe_name(Path(name).stem)}_highlight.doc"
        artifact_type, strategy = _highlight_doc(source, target, hit_values)
        artifact_path = target if artifact_type == "doc_highlight" else target.with_suffix(".docx")
        return (str(artifact_path) if artifact_type else None, artifact_type, {"status": "ok" if artifact_type else strategy, "strategy": strategy})
    if suffix == ".pdf":
        target = target_dir / f"v{version_no}_{_safe_name(Path(name).stem)}_native_highlight.pdf"
        pdf_info = _classify_pdf_file(source)
        artifact_type, applied = _highlight_pdf(source, target, hit_values)
        if artifact_type:
            return (str(target), artifact_type, {"status": "ok", "strategy": "pdf_native_textmarkup_annotation", "applied_count": applied, **pdf_info})
        target = target_dir / f"v{version_no}_{_safe_name(Path(name).stem)}_ocr_bbox_highlight.pdf"
        artifact_type, applied = _highlight_pdf_ocr_boxes(source, target, hits)
        status = "ok" if artifact_type else "no_pdf_coordinates"
        return (str(target) if artifact_type else None, artifact_type, {"status": status, "strategy": "pdf_ocr_bbox_annotation", "applied_count": applied, **pdf_info})
    if suffix == ".pptx":
        target = target_dir / f"v{version_no}_{_safe_name(Path(name).stem)}_highlight.pptx"
        artifact_type, applied = _highlight_pptx(source, target, hit_values)
        return (str(target) if artifact_type else None, artifact_type, {"status": "ok" if artifact_type else "no_hits", "strategy": "pptx_shape_or_cell_fill", "applied_count": applied})
    if suffix == ".ppt":
        target = target_dir / f"v{version_no}_{_safe_name(Path(name).stem)}_highlight.ppt"
        artifact_type, applied, strategy = _highlight_ppt(source, target, hit_values)
        artifact_path = target if artifact_type == "ppt_shape_highlight" else target.with_suffix(".pptx")
        return (str(artifact_path) if artifact_type else None, artifact_type, {"status": "ok" if artifact_type else strategy, "strategy": strategy, "applied_count": applied})
    if suffix == ".xlsx":
        target = target_dir / f"v{version_no}_{_safe_name(Path(name).stem)}_highlight.xlsx"
        artifact_type, applied = _highlight_xlsx(source, target, hit_values)
        return (str(target) if artifact_type else None, artifact_type, {"status": "ok" if artifact_type else "no_hits", "strategy": "xlsx_cell_fill", "applied_count": applied})
    if suffix == ".csv":
        target = target_dir / f"v{version_no}_{_safe_name(Path(name).stem)}_highlight.xlsx"
        artifact_type, applied = _highlight_csv(source, target, hit_values)
        return (str(target) if artifact_type else None, artifact_type, {"status": "ok" if artifact_type else "no_hits", "strategy": "csv_export_to_xlsx_cell_fill", "applied_count": applied})
    return None, None, {"status": "unsupported_file_type"}


def _copy_snapshot(file_row: FileRecord, dirs: dict[str, Path], version_no: int, name: str) -> Optional[Path]:
    if not file_row.store_path:
        return None
    suffix = Path(name or file_row.file_name or "").suffix or (file_row.file_type or "")
    target = dirs["versions"] / f"v{version_no}_{_safe_name(name or file_row.file_name or file_row.file_hash)}{suffix if suffix and not str(name).lower().endswith(str(suffix).lower()) else ''}"
    data = object_storage.get_bytes(file_row.store_path)
    target.write_bytes(data)
    return target


def _find_tracked(session, agent_id: str, path: Optional[str] = None, content_hash: Optional[str] = None) -> Optional[TrackedFile]:
    query = session.query(TrackedFile).filter(TrackedFile.agent_id == agent_id)
    file_key = _file_key(agent_id, path or "") if path else ""
    if file_key:
        row = query.filter(TrackedFile.file_key == file_key).first()
        if row:
            return row
    if path:
        row = query.filter(TrackedFile.current_path == path).first()
        if row:
            return row
        norm = _norm_path(path)
        for item in query.all():
            if _norm_path(item.current_path) == norm or _norm_path(item.original_path) == norm:
                return item
    if content_hash:
        version = session.query(TrackedFileVersion).filter(TrackedFileVersion.content_hash == content_hash).order_by(TrackedFileVersion.snapshot_time.desc()).first()
        if version:
            return session.get(TrackedFile, version.tracked_file_id)
    return None


def _serialize_tracked_file(row: TrackedFile, latest_event_type: str = "") -> dict:
    return {
        "tracked_file_id": row.tracked_file_id,
        "agent_id": row.agent_id,
        "file_key": row.file_key,
        "current_path": row.current_path,
        "current_name": row.current_name,
        "original_path": row.original_path,
        "original_name": row.original_name,
        "file_type": row.file_type,
        "sensitive_level": row.sensitive_level,
        "is_deleted": row.is_deleted,
        "first_seen_at": row.first_seen_at,
        "last_seen_at": row.last_seen_at,
        "deleted_at": row.deleted_at,
        "latest_event_type": latest_event_type,
        "latest_version_no": row.latest_version_no,
        "latest_version_id": row.latest_version_id,
        "rename_count": row.rename_count,
        "modify_count": row.modify_count,
    }


def _serialize_version(row: TrackedFileVersion, tracked_file_id: Optional[str] = None) -> dict:
    base = f"/api/v1/sensitive-files/{tracked_file_id or row.tracked_file_id}/versions/{row.version_id}"
    return {
        "version_id": row.version_id,
        "tracked_file_id": row.tracked_file_id,
        "version_no": row.version_no,
        "snapshot_time": row.snapshot_time,
        "event_type": row.event_type,
        "path_at_that_time": row.path_at_that_time,
        "name_at_that_time": row.name_at_that_time,
        "content_hash": row.content_hash,
        "prev_version_id": row.prev_version_id,
        "change_summary": row.change_summary,
        "change_detail_json": row.change_detail_json or {},
        "sensitive_hits": row.sensitive_hits or [],
        "can_download": bool(row.stored_file_path and row.is_snapshot_retained and Path(row.stored_file_path).exists()),
        "has_highlight": bool(row.highlight_artifact_path and Path(row.highlight_artifact_path).exists()),
        "has_diff": bool(int(row.version_no or 0) > 1 and row.diff_artifact_path and Path(row.diff_artifact_path).exists()),
        "download_url": f"{base}/download",
        "highlight_download_url": f"{base}/download-highlight",
        "diff_download_url": f"{base}/download-diff",
        "snapshot_retention_note": "" if row.is_snapshot_retained else "该历史版本仅保留摘要，原文件快照已清理",
        "artifact_type": row.artifact_type,
    }


def _serialize_event(row: TrackedFileEvent) -> dict:
    return {
        "event_id": row.event_id,
        "tracked_file_id": row.tracked_file_id,
        "event_time": row.event_time,
        "event_type": row.event_type,
        "old_path": row.old_path,
        "new_path": row.new_path,
        "old_name": row.old_name,
        "new_name": row.new_name,
        "description": row.description,
        "raw_event_json": row.raw_event_json or {},
        "version_id": row.version_id,
    }


def _path_history_from_events(tracked: TrackedFile, events: list[TrackedFileEvent]) -> list[dict[str, Any]]:
    history = []
    for event in events:
        if event.event_type not in {"modified", "renamed", "deleted", "restored"}:
            continue
        if event.old_path == event.new_path and event.event_type not in {"deleted", "renamed"}:
            continue
        history.append(
            {
                "event_time": event.event_time,
                "event_type": event.event_type,
                "old_path": event.old_path,
                "new_path": event.new_path,
                "old_name": event.old_name,
                "new_name": event.new_name,
                "version_id": event.version_id,
            }
        )
    deduped = []
    seen = set()
    for item in sorted(history, key=lambda value: float(value.get("event_time") or 0)):
        key = (item.get("event_type"), item.get("old_path"), item.get("new_path"), item.get("version_id"))
        if key in seen:
            continue
        seen.add(key)
        deduped.append(item)
    return deduped


def _sync_metadata(session, tracked: TrackedFile) -> None:
    dirs = _ensure_dirs(tracked.agent_id, tracked.tracked_file_id)
    latest_event = session.query(TrackedFileEvent).filter(TrackedFileEvent.tracked_file_id == tracked.tracked_file_id).order_by(TrackedFileEvent.event_time.desc()).first()
    _write_json(dirs["root"] / "metadata.json", _serialize_tracked_file(tracked, latest_event.event_type if latest_event else ""))


def _merge_duplicate_tracked_files(session) -> int:
    rows = session.query(TrackedFile).order_by(TrackedFile.first_seen_at.asc(), TrackedFile.last_seen_at.asc()).all()
    by_key: dict[tuple[str, str], TrackedFile] = {}
    merged = 0
    for row in rows:
        stable_key = _file_key(row.agent_id, row.current_path or row.original_path or "")
        row.file_key = stable_key
        key = (row.agent_id, stable_key)
        keeper = by_key.get(key)
        if not keeper:
            by_key[key] = row
            continue
        if float(row.last_seen_at or 0) > float(keeper.last_seen_at or 0):
            keeper.current_path = row.current_path
            keeper.current_name = row.current_name
            keeper.last_seen_at = row.last_seen_at
            keeper.latest_version_no = max(int(keeper.latest_version_no or 0), int(row.latest_version_no or 0))
            keeper.latest_version_id = row.latest_version_id or keeper.latest_version_id
        keeper.first_seen_at = min(float(keeper.first_seen_at or row.first_seen_at or _now()), float(row.first_seen_at or _now()))
        keeper.rename_count = int(keeper.rename_count or 0) + int(row.rename_count or 0)
        keeper.modify_count = int(keeper.modify_count or 0) + int(row.modify_count or 0)
        keeper.is_deleted = bool(keeper.is_deleted and row.is_deleted)
        keeper.deleted_at = keeper.deleted_at if keeper.is_deleted else None

        versions = (
            session.query(TrackedFileVersion)
            .filter(TrackedFileVersion.tracked_file_id == row.tracked_file_id)
            .order_by(TrackedFileVersion.version_no.asc(), TrackedFileVersion.snapshot_time.asc())
            .all()
        )
        next_version_no = int(keeper.latest_version_no or 0)
        for version in versions:
            existing = None
            if version.content_hash:
                existing = (
                    session.query(TrackedFileVersion)
                    .filter(
                        TrackedFileVersion.tracked_file_id == keeper.tracked_file_id,
                        TrackedFileVersion.content_hash == version.content_hash,
                        TrackedFileVersion.path_at_that_time == version.path_at_that_time,
                    )
                    .first()
                )
            if existing:
                session.query(TrackedFileEvent).filter(TrackedFileEvent.version_id == version.version_id).delete(synchronize_session=False)
                session.delete(version)
                continue
            next_version_no += 1
            version.tracked_file_id = keeper.tracked_file_id
            version.version_no = next_version_no
            session.add(version)
            keeper.latest_version_id = version.version_id
            keeper.latest_version_no = next_version_no
        session.query(TrackedFileEvent).filter(TrackedFileEvent.tracked_file_id == row.tracked_file_id).update(
            {"tracked_file_id": keeper.tracked_file_id},
            synchronize_session=False,
        )
        session.delete(row)
        merged += 1
    if merged:
        session.flush()
        for tracked in by_key.values():
            _sync_metadata(session, tracked)
    return merged


def _dedupe_tracked_versions(session, tracked: TrackedFile) -> int:
    versions = (
        session.query(TrackedFileVersion)
        .filter(TrackedFileVersion.tracked_file_id == tracked.tracked_file_id)
        .order_by(TrackedFileVersion.snapshot_time.asc(), TrackedFileVersion.version_no.asc())
        .all()
    )
    removed_ids = set()
    by_id = {version.version_id: version for version in versions}
    for version in versions:
        if version.event_type != "renamed":
            continue
        detail = dict(version.change_detail_json or {})
        rename = dict(detail.get("rename") or {})
        old_path = rename.get("old_path") or ""
        new_path = rename.get("new_path") or version.path_at_that_time
        if not _is_download_temp_path(old_path):
            continue
        prev = by_id.get(version.prev_version_id or "")
        stable_old_path = ""
        if prev and prev.path_at_that_time and not _is_download_temp_path(prev.path_at_that_time):
            stable_old_path = prev.path_at_that_time
        elif tracked.original_path and not _is_download_temp_path(tracked.original_path):
            stable_old_path = tracked.original_path
        if not stable_old_path or _norm_path(stable_old_path) == _norm_path(new_path):
            session.query(TrackedFileEvent).filter(TrackedFileEvent.version_id == version.version_id).delete(synchronize_session=False)
            session.delete(version)
            removed_ids.add(version.version_id)
            continue
        old_name = remote_path_name(stable_old_path, rename.get("old_name") or tracked.original_name or "file")
        new_name = remote_path_name(new_path, rename.get("new_name") or version.name_at_that_time or "file")
        rename.update({"old_path": stable_old_path, "new_path": new_path, "old_name": old_name, "new_name": new_name})
        detail["rename"] = rename
        detail["summary"] = f"文件已重命名：{old_name} -> {new_name}"
        version.change_detail_json = detail
        version.change_summary = detail["summary"]
        version.path_at_that_time = new_path
        version.name_at_that_time = new_name
        session.query(TrackedFileEvent).filter(TrackedFileEvent.version_id == version.version_id).update(
            {
                "old_path": stable_old_path,
                "new_path": new_path,
                "old_name": old_name,
                "new_name": new_name,
                "description": detail["summary"],
            },
            synchronize_session=False,
        )
        session.add(version)

    for version in versions:
        if version.version_id in removed_ids:
            continue
        hits = version.sensitive_hits or []
        detail = dict(version.change_detail_json or {})
        if not hits or detail.get("detection_summary"):
            continue
        detection_summary = _summarize_detection_hits(hits, {})
        if not detection_summary.get("summary"):
            continue
        detail["detection_summary"] = detection_summary
        detail["summary"] = _merge_summary_parts(detail.get("summary"), detection_summary.get("summary"))
        version.change_detail_json = detail
        if version.event_type == "initial":
            version.change_summary = _merge_summary_parts("首次敏感文件入库", detection_summary.get("summary"))
        else:
            version.change_summary = _merge_summary_parts(version.change_summary, detection_summary.get("summary"))
        session.add(version)

    seen = set()
    kept: list[TrackedFileVersion] = []
    removed = len(removed_ids)
    for version in versions:
        if version.version_id in removed_ids:
            continue
        key = (
            version.event_type,
            version.content_hash or "",
            _norm_path(version.path_at_that_time),
        )
        if version.event_type == "initial" and key in seen:
            session.query(TrackedFileEvent).filter(TrackedFileEvent.version_id == version.version_id).delete(synchronize_session=False)
            session.delete(version)
            removed += 1
            continue
        seen.add(key)
        kept.append(version)
    if removed:
        for idx, version in enumerate(kept, start=1):
            version.version_no = idx
            session.add(version)
        latest = kept[-1] if kept else None
        tracked.latest_version_no = int(latest.version_no) if latest else 0
        tracked.latest_version_id = latest.version_id if latest else None
        session.flush()
        _sync_metadata(session, tracked)
    return removed


def _prune_old_snapshots(session, tracked: TrackedFile) -> None:
    rows = (
        session.query(TrackedFileVersion)
        .filter(TrackedFileVersion.tracked_file_id == tracked.tracked_file_id, TrackedFileVersion.event_type.in_(list(CONTENT_EVENTS)))
        .order_by(TrackedFileVersion.version_no.desc())
        .all()
    )
    for row in rows[RETAINED_CONTENT_SNAPSHOTS:]:
        if not row.stored_file_path or not row.is_snapshot_retained:
            continue
        path = Path(row.stored_file_path)
        try:
            if path.exists():
                path.unlink()
        except Exception:
            pass
        row.is_snapshot_retained = False
        row.stored_file_path = None
        session.add(row)


def archive_sensitive_file(file_hash: str, agent_id: Optional[str] = None, file_path: Optional[str] = None, event_type: str = "initial", event_time: Optional[float] = None) -> Optional[dict]:
    with db_session() as session:
        file_row = session.get(FileRecord, file_hash)
        parse_row = session.get(ParseResult, file_hash)
        if not file_row or not file_row.is_sensitive:
            return None
        versions_query = session.query(FileVersion).filter(FileVersion.file_hash == file_hash)
        if agent_id:
            versions_query = versions_query.filter(FileVersion.agent_id == agent_id)
        if file_path:
            versions_query = versions_query.filter(FileVersion.file_path == file_path)
        file_versions = versions_query.order_by(FileVersion.is_current.desc(), FileVersion.created_at.asc()).all()
        if not file_versions:
            file_versions = [FileVersion(agent_id=agent_id or "server", file_path=file_path or file_row.file_name or file_hash, file_hash=file_hash, is_current=True)]

        deduped_versions = []
        seen_keys = set()
        for current in file_versions:
            key = (str(current.agent_id or ""), _file_key(str(current.agent_id or ""), str(current.file_path or "")))
            if key in seen_keys:
                continue
            seen_keys.add(key)
            deduped_versions.append(current)

        archived = None
        for current in deduped_versions:
            archived = _archive_sensitive_row(session, file_row, parse_row, current.agent_id, current.file_path, event_type, event_time)
        return archived


def _archive_sensitive_row(session, file_row: FileRecord, parse_row: Optional[ParseResult], agent_id: str, path: str, event_type: str, event_time: Optional[float]) -> dict:
    event_time = float(event_time or _now())
    name = remote_path_name(path or file_row.file_name, file_row.file_hash)
    content_hash_lookup = None if event_type == "initial" else file_row.file_hash
    tracked = _find_tracked(session, agent_id, path=path, content_hash=content_hash_lookup)
    if not tracked:
        tracked = TrackedFile(
            tracked_file_id=str(uuid.uuid4()),
            agent_id=agent_id,
            file_key=_file_key(agent_id, path),
            current_path=path,
            current_name=name,
            original_path=path,
            original_name=name,
            file_type=file_row.file_type or Path(name).suffix.lower(),
            sensitive_level=file_row.risk_level,
            is_deleted=False,
            first_seen_at=event_time,
            last_seen_at=event_time,
            latest_version_no=0,
            rename_count=0,
            modify_count=0,
        )
        session.add(tracked)
        session.flush()

    existing_hash = (
        session.query(TrackedFileVersion)
        .filter(TrackedFileVersion.tracked_file_id == tracked.tracked_file_id, TrackedFileVersion.content_hash == file_row.file_hash)
        .order_by(TrackedFileVersion.version_no.desc())
        .first()
    )
    if existing_hash and event_type == "initial":
        tracked.current_path = path
        tracked.current_name = name
        tracked.file_key = _file_key(agent_id, path)
        tracked.last_seen_at = event_time
        _sync_metadata(session, tracked)
        return _serialize_tracked_file(tracked)

    prev = session.get(TrackedFileVersion, tracked.latest_version_id) if tracked.latest_version_id else None
    actual_event = "initial" if tracked.latest_version_no <= 0 else ("modified" if event_type == "initial" else event_type)
    version_no = int(tracked.latest_version_no or 0) + 1
    dirs = _ensure_dirs(agent_id, tracked.tracked_file_id)
    snapshot_path = _copy_snapshot(file_row, dirs, version_no, name)
    old_path = Path(prev.stored_file_path) if prev and prev.stored_file_path and Path(prev.stored_file_path).exists() else None
    diff = _diff_texts(old_path, snapshot_path)
    diff_path = dirs["diffs"] / f"v{version_no}_diff.json"
    _write_json(diff_path, diff)
    parse_data = parse_row.result_data if parse_row else {}
    hits = _extract_hits(parse_data)
    detection_summary = _summarize_detection_hits(hits, parse_data)
    if detection_summary.get("summary"):
        diff["detection_summary"] = detection_summary
        diff["summary"] = _merge_summary_parts(diff.get("summary"), detection_summary.get("summary"))
    highlight_path, artifact_type, highlight_detail = _create_highlight(snapshot_path, dirs["highlights"], version_no, name, hits)
    diff["highlight"] = highlight_detail
    _write_json(diff_path, diff)
    initial_summary = _merge_summary_parts("首次敏感文件入库", detection_summary.get("summary"))

    version = TrackedFileVersion(
        version_id=str(uuid.uuid4()),
        tracked_file_id=tracked.tracked_file_id,
        version_no=version_no,
        snapshot_time=event_time,
        event_type=actual_event,
        path_at_that_time=path,
        name_at_that_time=name,
        stored_file_path=str(snapshot_path) if snapshot_path else None,
        content_hash=file_row.file_hash,
        prev_version_id=tracked.latest_version_id,
        change_summary=initial_summary if actual_event == "initial" else diff["summary"],
        change_detail_json=diff,
        highlight_artifact_path=highlight_path,
        diff_artifact_path=str(diff_path),
        sensitive_hits=hits,
        artifact_type=artifact_type,
        is_snapshot_retained=bool(snapshot_path),
    )
    session.add(version)
    tracked.current_path = path
    tracked.current_name = name
    tracked.file_key = _file_key(agent_id, path)
    tracked.file_type = file_row.file_type or tracked.file_type
    tracked.sensitive_level = file_row.risk_level or tracked.sensitive_level
    tracked.is_deleted = False
    tracked.deleted_at = None
    tracked.last_seen_at = event_time
    tracked.latest_version_no = version_no
    tracked.latest_version_id = version.version_id
    if actual_event == "modified":
        tracked.modify_count = int(tracked.modify_count or 0) + 1
    event = TrackedFileEvent(
        event_id=f"version:{version.version_id}",
        tracked_file_id=tracked.tracked_file_id,
        event_time=event_time,
        event_type=actual_event,
        old_path=prev.path_at_that_time if prev else None,
        new_path=path,
        old_name=prev.name_at_that_time if prev else None,
        new_name=name,
        description=version.change_summary,
        raw_event_json={"file_hash": file_row.file_hash},
        version_id=version.version_id,
    )
    session.add(event)
    session.flush()
    _prune_old_snapshots(session, tracked)
    _sync_metadata(session, tracked)
    _write_json(dirs["events"] / f"{event.event_id.replace(':', '_')}.json", _serialize_event(event))
    return _serialize_tracked_file(tracked, actual_event)


def ingest_tracked_event(item: dict) -> Optional[dict]:
    event_type = str(item.get("event_type") or "").lower()
    if event_type in {"file_renamed", "renamed", "rename", "file_moved"}:
        return record_rename_event(item)
    if event_type in {"file_deleted", "deleted", "delete"}:
        return record_delete_event(item)
    if event_type in {"file_modified", "file_overwritten", "modified", "file_changed"}:
        file_hash = item.get("new_hash") or item.get("old_hash")
        if file_hash:
            return archive_sensitive_file(
                str(file_hash),
                agent_id=item.get("agent_id"),
                file_path=item.get("new_path") or item.get("file_path") or item.get("old_path"),
                event_type="modified",
                event_time=float(item.get("timestamp") or _now()),
            )
    return None


def record_rename_event(item: dict) -> Optional[dict]:
    with db_session() as session:
        agent_id = item.get("agent_id")
        old_path = item.get("old_path") or item.get("file_path")
        new_path = item.get("new_path") or item.get("file_path")
        tracked = _find_tracked(session, agent_id, path=old_path, content_hash=item.get("old_hash") or item.get("new_hash"))
        if not tracked or not new_path:
            return None
        if _is_download_temp_path(old_path):
            stable_old_path = tracked.current_path or tracked.original_path
            if _norm_path(stable_old_path) == _norm_path(new_path):
                return None
            if stable_old_path and not _is_download_temp_path(stable_old_path):
                old_path = stable_old_path
        event_time = float(item.get("timestamp") or _now())
        old_name = remote_path_name(old_path or tracked.current_path, tracked.current_name or "file")
        new_name = remote_path_name(new_path, tracked.current_name or "file")
        prev = session.get(TrackedFileVersion, tracked.latest_version_id) if tracked.latest_version_id else None
        version_no = int(tracked.latest_version_no or 0) + 1
        detail = {
            "summary": f"文件已重命名：{old_name} -> {new_name}",
            "added_texts": [],
            "removed_texts": [],
            "modified_blocks": [],
            "rename": {"old_name": old_name, "new_name": new_name, "old_path": old_path, "new_path": new_path},
            "deleted": None,
        }
        dirs = _ensure_dirs(agent_id, tracked.tracked_file_id)
        diff_path = dirs["diffs"] / f"v{version_no}_rename.json"
        _write_json(diff_path, detail)
        version = TrackedFileVersion(
            version_id=str(uuid.uuid4()),
            tracked_file_id=tracked.tracked_file_id,
            version_no=version_no,
            snapshot_time=event_time,
            event_type="renamed",
            path_at_that_time=new_path,
            name_at_that_time=new_name,
            stored_file_path=prev.stored_file_path if prev else None,
            content_hash=prev.content_hash if prev else (item.get("new_hash") or item.get("old_hash")),
            prev_version_id=tracked.latest_version_id,
            change_summary=detail["summary"],
            change_detail_json=detail,
            diff_artifact_path=str(diff_path),
            sensitive_hits=prev.sensitive_hits if prev else [],
            is_snapshot_retained=bool(prev and prev.stored_file_path and Path(prev.stored_file_path).exists()),
            artifact_type=prev.artifact_type if prev else None,
            highlight_artifact_path=prev.highlight_artifact_path if prev else None,
        )
        event = TrackedFileEvent(
            event_id=item.get("event_id") or str(uuid.uuid4()),
            tracked_file_id=tracked.tracked_file_id,
            event_time=event_time,
            event_type="renamed",
            old_path=old_path,
            new_path=new_path,
            old_name=old_name,
            new_name=new_name,
            description=detail["summary"],
            raw_event_json=item,
            version_id=version.version_id,
        )
        session.add(version)
        session.add(event)
        tracked.current_path = new_path
        tracked.current_name = new_name
        tracked.file_key = _file_key(agent_id, new_path)
        tracked.last_seen_at = event_time
        tracked.latest_version_no = version_no
        tracked.latest_version_id = version.version_id
        tracked.rename_count = int(tracked.rename_count or 0) + 1
        session.flush()
        _sync_metadata(session, tracked)
        _write_json(dirs["events"] / f"{event.event_id}.json", _serialize_event(event))
        return _serialize_tracked_file(tracked, "renamed")


def record_delete_event(item: dict) -> Optional[dict]:
    with db_session() as session:
        agent_id = item.get("agent_id")
        path = item.get("file_path") or item.get("old_path") or item.get("new_path")
        tracked = _find_tracked(session, agent_id, path=path, content_hash=item.get("old_hash") or item.get("new_hash"))
        if not tracked:
            return None
        event_time = float(item.get("timestamp") or _now())
        prev = session.get(TrackedFileVersion, tracked.latest_version_id) if tracked.latest_version_id else None
        version_no = int(tracked.latest_version_no or 0) + 1
        deleted = {"is_deleted": True, "deleted_at": event_time}
        detail = _diff_texts(None, None, deleted=deleted)
        dirs = _ensure_dirs(agent_id, tracked.tracked_file_id)
        diff_path = dirs["diffs"] / f"v{version_no}_deleted.json"
        _write_json(diff_path, detail)
        version = TrackedFileVersion(
            version_id=str(uuid.uuid4()),
            tracked_file_id=tracked.tracked_file_id,
            version_no=version_no,
            snapshot_time=event_time,
            event_type="deleted",
            path_at_that_time=tracked.current_path,
            name_at_that_time=tracked.current_name,
            stored_file_path=prev.stored_file_path if prev else None,
            content_hash=prev.content_hash if prev else (item.get("old_hash") or item.get("new_hash")),
            prev_version_id=tracked.latest_version_id,
            change_summary=detail["summary"],
            change_detail_json=detail,
            diff_artifact_path=str(diff_path),
            sensitive_hits=prev.sensitive_hits if prev else [],
            is_snapshot_retained=bool(prev and prev.stored_file_path and Path(prev.stored_file_path).exists()),
            artifact_type=prev.artifact_type if prev else None,
            highlight_artifact_path=prev.highlight_artifact_path if prev else None,
        )
        event = TrackedFileEvent(
            event_id=item.get("event_id") or str(uuid.uuid4()),
            tracked_file_id=tracked.tracked_file_id,
            event_time=event_time,
            event_type="deleted",
            old_path=tracked.current_path,
            new_path=None,
            old_name=tracked.current_name,
            new_name=None,
            description=detail["summary"],
            raw_event_json=item,
            version_id=version.version_id,
        )
        session.add(version)
        session.add(event)
        tracked.is_deleted = True
        tracked.deleted_at = event_time
        tracked.last_seen_at = event_time
        tracked.latest_version_no = version_no
        tracked.latest_version_id = version.version_id
        session.flush()
        _sync_metadata(session, tracked)
        _write_json(dirs["events"] / f"{event.event_id}.json", _serialize_event(event))
        return _serialize_tracked_file(tracked, "deleted")


def list_sensitive_files(agent_id: Optional[str] = None, changed_only: bool = False, is_deleted: Optional[bool] = None, keyword: Optional[str] = None, file_type: Optional[str] = None, page: int = 1, page_size: int = 50) -> dict:
    backfill_sensitive_archives(agent_id=agent_id)
    page = max(1, int(page or 1))
    page_size = min(200, max(1, int(page_size or 50)))
    with db_session() as session:
        _merge_duplicate_tracked_files(session)
        query = session.query(TrackedFile)
        if agent_id:
            query = query.filter(TrackedFile.agent_id == agent_id)
        if is_deleted is not None:
            query = query.filter(TrackedFile.is_deleted.is_(bool(is_deleted)))
        if changed_only:
            query = query.filter((TrackedFile.modify_count > 0) | (TrackedFile.rename_count > 0) | (TrackedFile.is_deleted.is_(True)))
        if file_type:
            query = query.filter(TrackedFile.file_type == file_type)
        rows = query.order_by(TrackedFile.last_seen_at.desc()).all()
        if keyword:
            needle = str(keyword).lower()
            rows = [row for row in rows if needle in (row.current_name or "").lower() or needle in (row.current_path or "").lower()]
        total = len(rows)
        rows = rows[(page - 1) * page_size : page * page_size]
        for row in rows:
            _dedupe_tracked_versions(session, row)
        latest_events = {}
        for event in session.query(TrackedFileEvent).filter(TrackedFileEvent.tracked_file_id.in_([row.tracked_file_id for row in rows] or [""])).order_by(TrackedFileEvent.event_time.desc()).all():
            if event.tracked_file_id and event.tracked_file_id not in latest_events:
                latest_events[event.tracked_file_id] = event.event_type
        return {
            "items": [_serialize_tracked_file(row, latest_events.get(row.tracked_file_id, "")) for row in rows],
            "page": page,
            "page_size": page_size,
            "total": total,
        }


def backfill_sensitive_archives(agent_id: Optional[str] = None, limit: int = 50) -> int:
    created = 0
    with db_session() as session:
        query = (
            session.query(FileVersion)
            .join(FileRecord, FileRecord.file_hash == FileVersion.file_hash)
            .filter(FileRecord.is_sensitive.is_(True))
            .order_by(FileVersion.created_at.asc())
        )
        if agent_id:
            query = query.filter(FileVersion.agent_id == agent_id)
        rows = query.limit(limit).all()
        for row in rows:
            if _find_tracked(session, row.agent_id, path=row.file_path):
                continue
            file_row = session.get(FileRecord, row.file_hash)
            parse_row = session.get(ParseResult, row.file_hash)
            if not file_row:
                continue
            _archive_sensitive_row(session, file_row, parse_row, row.agent_id, row.file_path, "initial", row.created_at)
            created += 1
    return created


def get_sensitive_file_history(tracked_file_id: str) -> dict:
    with db_session() as session:
        _merge_duplicate_tracked_files(session)
        tracked = session.get(TrackedFile, tracked_file_id)
        if not tracked:
            raise ValueError("tracked file not found")
        _dedupe_tracked_versions(session, tracked)
        versions = session.query(TrackedFileVersion).filter(TrackedFileVersion.tracked_file_id == tracked_file_id).order_by(TrackedFileVersion.snapshot_time.asc(), TrackedFileVersion.version_no.asc()).all()
        events = session.query(TrackedFileEvent).filter(TrackedFileEvent.tracked_file_id == tracked_file_id).order_by(TrackedFileEvent.event_time.asc()).all()
        latest_event = events[-1].event_type if events else ""
        return {
            "file": _serialize_tracked_file(tracked, latest_event),
            "current_state": "deleted" if tracked.is_deleted else "active",
            "versions": [_serialize_version(row, tracked_file_id) for row in versions],
            "events": [_serialize_event(row) for row in events],
            "path_history": _path_history_from_events(tracked, events),
        }


def get_sensitive_version_detail(tracked_file_id: str, version_id: str) -> dict:
    with db_session() as session:
        tracked = session.get(TrackedFile, tracked_file_id)
        version = session.get(TrackedFileVersion, version_id)
        if not tracked or not version or version.tracked_file_id != tracked_file_id:
            raise ValueError("version not found")
        payload = _serialize_version(version, tracked_file_id)
        payload.update(
            {
                "file": _serialize_tracked_file(tracked),
                "change_summary": version.change_summary,
                "change_detail_json": version.change_detail_json or {},
                "sensitive_hits": version.sensitive_hits or [],
            }
        )
        return payload


def get_version_artifact_path(tracked_file_id: str, version_id: str, artifact: str) -> Path:
    info = get_version_artifact_info(tracked_file_id, version_id, artifact)
    return info["path"]


def _download_name_for_version(version: TrackedFileVersion, artifact: str, path: Path) -> str:
    original = Path(version.name_at_that_time or version.path_at_that_time or "archive").name
    suffix = path.suffix or Path(original).suffix
    stem = Path(original).stem or "archive"
    stem = re.sub(r'[<>:"/\\|?*]+', "_", stem).strip() or "archive"
    version_label = f"V{int(version.version_no or 0)}"
    if artifact == "source":
        return f"{stem}_{version_label}{suffix}"
    if artifact == "highlight":
        return f"{stem}_{version_label}_高亮{suffix}"
    return f"{stem}_{version_label}_差异.json"


def get_version_artifact_info(tracked_file_id: str, version_id: str, artifact: str) -> dict:
    with db_session() as session:
        version = session.get(TrackedFileVersion, version_id)
        if not version or version.tracked_file_id != tracked_file_id:
            raise ValueError("version not found")
        if artifact == "source":
            if not version.stored_file_path or not version.is_snapshot_retained:
                raise FileNotFoundError("snapshot file has been cleaned")
            path = Path(version.stored_file_path)
        elif artifact == "highlight":
            if not version.highlight_artifact_path:
                raise FileNotFoundError("highlight artifact not found")
            path = Path(version.highlight_artifact_path)
        elif artifact == "diff":
            if int(version.version_no or 0) <= 1:
                raise FileNotFoundError("initial version has no diff artifact")
            if not version.diff_artifact_path:
                raise FileNotFoundError("diff artifact not found")
            path = Path(version.diff_artifact_path)
        else:
            raise ValueError("unsupported artifact")
        root = TRACK_ROOT.resolve()
        resolved = path.resolve()
        if root not in resolved.parents:
            raise ValueError("artifact path outside guard_state")
        if not resolved.exists():
            raise FileNotFoundError("artifact file not found")
        return {"path": resolved, "filename": _download_name_for_version(version, artifact, resolved)}
