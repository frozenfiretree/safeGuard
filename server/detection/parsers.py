import csv
import io
import shutil
import subprocess
import tempfile
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Dict, List, Optional


SUPPORTED_EXTENSIONS = {
    ".txt", ".csv", ".doc", ".docx", ".pdf", ".xlsx", ".pptx",
    ".png", ".jpg", ".jpeg", ".bmp",
}

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".bmp"}


def make_text_block(text: str, location: str, source_type: str) -> Dict:
    return {
        "text": text,
        "location": location,
        "source_type": source_type,
    }


def make_image_block(image_bytes: bytes, location: str, source_type: str) -> Dict:
    return {
        "bytes": image_bytes,
        "location": location,
        "source_type": source_type,
    }


def _extract_txt(path: Path) -> Dict:
    blocks = []
    for idx, line in enumerate(path.read_text(encoding="utf-8", errors="ignore").splitlines(), start=1):
        line = line.strip()
        if line:
            blocks.append(make_text_block(line, f"line:{idx}", "text_line"))
    return {"text_blocks": blocks, "image_blocks": [], "needs_ocr": False, "parse_status": "ok"}


def _extract_csv(path: Path) -> Dict:
    blocks = []
    with open(path, "r", encoding="utf-8", errors="ignore", newline="") as f:
        reader = csv.reader(f)
        for row_idx, row in enumerate(reader, start=1):
            for col_idx, value in enumerate(row, start=1):
                text = str(value).strip()
                if text:
                    blocks.append(make_text_block(text, f"row:{row_idx}:col:{col_idx}", "csv_cell"))
    return {"text_blocks": blocks, "image_blocks": [], "needs_ocr": False, "parse_status": "ok"}


def _extract_docx(path: Path) -> Dict:
    from docx import Document

    doc = Document(path)
    text_blocks = []
    image_blocks = []

    for idx, para in enumerate(doc.paragraphs, start=1):
        text = (para.text or "").strip()
        if text:
            text_blocks.append(make_text_block(text, f"paragraph:{idx}", "docx_paragraph"))

    for table_idx, table in enumerate(doc.tables, start=1):
        for row_idx, row in enumerate(table.rows, start=1):
            for col_idx, cell in enumerate(row.cells, start=1):
                text = (cell.text or "").strip()
                if text:
                    text_blocks.append(
                        make_text_block(text, f"table:{table_idx}:r{row_idx}:c{col_idx}", "docx_table_cell")
                    )

    rels = getattr(doc.part, "rels", {})
    image_index = 0
    for rel in rels.values():
        target = getattr(rel, "target_ref", "")
        if "image" not in target:
            continue
        image_part = getattr(rel, "target_part", None)
        blob = getattr(image_part, "blob", None)
        if blob:
            image_index += 1
            image_blocks.append(make_image_block(blob, f"docx:image:{image_index}", "docx_image"))

    return {
        "text_blocks": text_blocks,
        "image_blocks": image_blocks,
        "needs_ocr": bool(image_blocks),
        "parse_status": "ok",
    }


def _extract_xlsx(path: Path) -> Dict:
    import openpyxl

    wb = openpyxl.load_workbook(path, data_only=True)
    text_blocks = []
    image_blocks = []

    for sheet in wb.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                value = cell.value
                if value is None:
                    continue
                text = str(value).strip()
                if text:
                    text_blocks.append(make_text_block(text, f"sheet:{sheet.title}:{cell.coordinate}", "xlsx_cell"))

    try:
        with zipfile.ZipFile(path, "r") as zf:
            image_index = 0
            for name in zf.namelist():
                if name.startswith("xl/media/"):
                    image_index += 1
                    image_blocks.append(
                        make_image_block(zf.read(name), f"xlsx:image:{image_index}:{Path(name).name}", "xlsx_image")
                    )
    except Exception:
        pass

    return {
        "text_blocks": text_blocks,
        "image_blocks": image_blocks,
        "needs_ocr": bool(image_blocks),
        "parse_status": "ok",
    }


def _extract_pptx(path: Path) -> Dict:
    from pptx import Presentation

    prs = Presentation(path)
    text_blocks = []
    image_blocks = []

    def walk_shapes(shapes, slide_idx: int, prefix: str = "shape"):
        nonlocal text_blocks
        nonlocal image_blocks
        for shape_idx, shape in enumerate(shapes, start=1):
            location = f"slide:{slide_idx}:{prefix}:{shape_idx}"
            if getattr(shape, "has_text_frame", False):
                text = (shape.text or "").strip()
                if text:
                    text_blocks.append(make_text_block(text, location, "pptx_shape"))
            if getattr(shape, "has_table", False):
                table = shape.table
                for row_idx, row in enumerate(table.rows, start=1):
                    for col_idx, cell in enumerate(row.cells, start=1):
                        text = (cell.text or "").strip()
                        if text:
                            text_blocks.append(make_text_block(text, f"{location}:table:r{row_idx}:c{col_idx}", "pptx_table_cell"))
            if hasattr(shape, "shapes"):
                walk_shapes(shape.shapes, slide_idx, f"{location}:group")
            try:
                blob = shape.image.blob
            except Exception:
                blob = None
            if blob:
                image_blocks.append(make_image_block(blob, f"{location}:image", "pptx_image"))

    for slide_idx, slide in enumerate(prs.slides, start=1):
        walk_shapes(slide.shapes, slide_idx)

    zip_texts, zip_images = _extract_pptx_from_zip(path)
    seen_texts = {(item["text"], item["location"]) for item in text_blocks}
    for item in zip_texts:
        key = (item["text"], item["location"])
        if key not in seen_texts:
            text_blocks.append(item)
            seen_texts.add(key)
    if not image_blocks:
        image_blocks.extend(zip_images)

    return {
        "text_blocks": text_blocks,
        "image_blocks": image_blocks,
        "needs_ocr": bool(image_blocks),
        "parse_status": "ok",
    }


def _extract_pptx_from_zip(path: Path) -> tuple[List[Dict], List[Dict]]:
    text_blocks: List[Dict] = []
    image_blocks: List[Dict] = []
    try:
        with zipfile.ZipFile(path, "r") as zf:
            slide_names = sorted(name for name in zf.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml"))
            ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
            for slide_idx, name in enumerate(slide_names, start=1):
                try:
                    root = ET.fromstring(zf.read(name))
                except Exception:
                    continue
                texts = [node.text.strip() for node in root.findall(".//a:t", ns) if node.text and node.text.strip()]
                if texts:
                    text_blocks.append(make_text_block("\n".join(texts), f"slide:{slide_idx}:xml", "pptx_xml_text"))
            image_index = 0
            for name in zf.namelist():
                lower = name.lower()
                if not lower.startswith("ppt/media/") or not lower.rsplit(".", 1)[-1] in {"png", "jpg", "jpeg", "bmp", "gif", "tif", "tiff"}:
                    continue
                image_index += 1
                image_blocks.append(make_image_block(zf.read(name), f"pptx:media:{image_index}:{Path(name).name}", "pptx_embedded_image"))
    except Exception:
        return [], []
    return text_blocks, image_blocks


def _extract_pdf(path: Path) -> Dict:
    import fitz

    doc = fitz.open(path)
    text_blocks = []
    image_blocks = []
    pages_without_text = 0

    for page_idx, page in enumerate(doc, start=1):
        text = (page.get_text("text") or "").strip()
        if text:
            for line_idx, line in enumerate(text.splitlines(), start=1):
                line = line.strip()
                if line:
                    text_blocks.append(make_text_block(line, f"page:{page_idx}:line:{line_idx}", "pdf_text"))
        else:
            pages_without_text += 1

        for img_idx, img in enumerate(page.get_images(full=True), start=1):
            try:
                xref = img[0]
                base = doc.extract_image(xref)
                blob = base.get("image")
                if blob:
                    image_blocks.append(
                        make_image_block(blob, f"page:{page_idx}:image:{img_idx}", "pdf_image")
                    )
            except Exception:
                continue

    return {
        "text_blocks": text_blocks,
        "image_blocks": image_blocks,
        "needs_ocr": bool(image_blocks) or pages_without_text > 0,
        "pages_without_text": pages_without_text,
        "parse_status": "ok",
    }


def _convert_doc_to_docx(path: Path) -> Optional[Path]:
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        return None

    temp_dir = Path(tempfile.mkdtemp(prefix="safeguard_doc_"))
    try:
        proc = subprocess.run(
            [
                soffice,
                "--headless",
                "--convert-to",
                "docx",
                "--outdir",
                str(temp_dir),
                str(path),
            ],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            errors="replace",
            timeout=120,
        )
        if proc.returncode != 0:
            shutil.rmtree(temp_dir, ignore_errors=True)
            return None
        out_path = temp_dir / f"{path.stem}.docx"
        if out_path.exists():
            return out_path
        shutil.rmtree(temp_dir, ignore_errors=True)
        return None
    except Exception:
        shutil.rmtree(temp_dir, ignore_errors=True)
        return None


def _extract_doc(path: Path) -> Dict:
    converted = _convert_doc_to_docx(path)
    if not converted:
        return {
            "text_blocks": [],
            "image_blocks": [],
            "needs_ocr": False,
            "parse_status": "unsupported_doc_conversion",
        }

    try:
        out = _extract_docx(converted)
        out["converted_from_doc"] = True
        return out
    finally:
        shutil.rmtree(converted.parent, ignore_errors=True)


def _extract_image(path: Path) -> Dict:
    from PIL import Image

    image = Image.open(path)
    buffer = io.BytesIO()
    image.save(buffer, format="PNG")
    return {
        "text_blocks": [],
        "image_blocks": [make_image_block(buffer.getvalue(), f"image:{path.name}", "image_file")],
        "needs_ocr": True,
        "parse_status": "ok",
    }


def extract_file_content(path: Path) -> Dict:
    ext = path.suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        return {"text_blocks": [], "image_blocks": [], "needs_ocr": False, "parse_status": "unsupported"}

    try:
        if ext == ".txt":
            return _extract_txt(path)
        if ext == ".csv":
            return _extract_csv(path)
        if ext == ".docx":
            return _extract_docx(path)
        if ext == ".xlsx":
            return _extract_xlsx(path)
        if ext == ".pptx":
            return _extract_pptx(path)
        if ext == ".pdf":
            return _extract_pdf(path)
        if ext == ".doc":
            return _extract_doc(path)
        if ext in IMAGE_EXTENSIONS:
            return _extract_image(path)
    except Exception as e:
        return {
            "text_blocks": [],
            "image_blocks": [],
            "needs_ocr": False,
            "parse_status": "failed",
            "error": str(e),
        }

    return {"text_blocks": [], "image_blocks": [], "needs_ocr": False, "parse_status": "unsupported"}
