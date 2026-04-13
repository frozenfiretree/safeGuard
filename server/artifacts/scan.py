import os
os.environ['MKLDNN_ENABLED'] = '0'
os.environ["FLAGS_use_mkldnn"] = "0"
os.environ["FLAGS_use_onednn"] = "0"
os.environ["FLAGS_enable_pir_api"] = "0"
os.environ["FLAGS_use_cinn"] = "0"
os.environ["FLAGS_use_gpu"] = "0"
os.environ["PADDLE_SKIP_LOAD_EXTENSION"] = "1"
os.environ["PADDLE_PDX_DISABLE_MODEL_SOURCE_CHECK"] = "True"
os.environ["GLOG_minloglevel"] = "2"
os.environ["FLAGS_minloglevel"] = "3"
os.environ["PADDLE_PDX_MODEL_SOURCE"] = "LOCAL"

import re
import io
import sys
import argparse
import string
import json
import traceback
import win32security
from tqdm import tqdm
import zipfile
import shutil
import tempfile
from datetime import datetime

import win32com.client
from docx import Document
from docx.enum.text import WD_COLOR_INDEX

import openpyxl
from openpyxl.styles import PatternFill

from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont

from paddleocr import PaddleOCR
import fitz
import numpy as np
import cv2

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as patches

# ===== matplotlib 中文字体设置，解决中文告警 =====
plt.rcParams['font.sans-serif'] = ['Microsoft YaHei', 'SimHei', 'DengXian', 'Arial Unicode MS']
plt.rcParams['axes.unicode_minus'] = False

OCR_MODEL = None
OCR_INIT_ERROR = None
TARGET_EXTENSIONS = ('.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx', '.txt', '.csv', '.pdf')
SKIP_FILE_PREFIXES = ('~$',)


# =========================================================
# 工具函数
# =========================================================

def find_matches_in_content(content, simple_keywords, regex_patterns):
    matches = []
    if not isinstance(content, str):
        return matches

    for keyword in simple_keywords:
        start = 0
        while (start := content.find(keyword, start)) != -1:
            matches.append(("keyword", keyword, start, start + len(keyword)))
            start += len(keyword)

    for pattern in regex_patterns:
        for match in re.finditer(pattern, content):
            matches.append(("regex", match.group(0), match.start(), match.end()))

    return matches


def normalize_box(box):
    if box is None:
        return None
    try:
        if isinstance(box, np.ndarray):
            box = box.tolist()

        if not isinstance(box, (list, tuple)) or len(box) == 0:
            return None

        if len(box) == 4 and not isinstance(box[0], (list, tuple)):
            x1, y1, x2, y2 = box
            return [
                [float(x1), float(y1)],
                [float(x2), float(y1)],
                [float(x2), float(y2)],
                [float(x1), float(y2)]
            ]

        if len(box) >= 4 and isinstance(box[0], (list, tuple, np.ndarray)):
            pts = []
            for p in box[:4]:
                if isinstance(p, np.ndarray):
                    p = p.tolist()
                if len(p) >= 2:
                    pts.append([float(p[0]), float(p[1])])
            if len(pts) == 4:
                return pts
    except Exception as e:
        print(f"[DEBUG] normalize_box 失败: {e}")
    return None


def image_bytes_to_cv2(img_bytes):
    try:
        nparr = np.frombuffer(img_bytes, np.uint8)
        return cv2.imdecode(nparr, cv2.IMREAD_COLOR)
    except Exception as e:
        print(f"[DEBUG] 图片解码失败: {e}")
        return None


def replace_bytes_in_zip(zip_file_path, replace_map):
    """
    将 zip_file_path 中的指定文件内容替换为 replace_map 中的字节，并回写原文件。
    为避免 WinError 17，临时文件创建在原文件同目录。
    """
    base_dir = os.path.dirname(os.path.abspath(zip_file_path))
    tmp_zip_path = os.path.join(base_dir, f"__tmp_replace_{os.path.basename(zip_file_path)}")

    try:
        with zipfile.ZipFile(zip_file_path, 'r') as zin:
            with zipfile.ZipFile(tmp_zip_path, 'w', zipfile.ZIP_DEFLATED) as zout:
                for item in zin.infolist():
                    data = zin.read(item.filename)
                    if item.filename in replace_map:
                        data = replace_map[item.filename]
                    zout.writestr(item, data)

        os.remove(zip_file_path)
        shutil.move(tmp_zip_path, zip_file_path)
        return True
    except Exception as e:
        print(f"[DEBUG] ZIP回写失败: {e}")
        try:
            if os.path.exists(tmp_zip_path):
                os.remove(tmp_zip_path)
        except:
            pass
        return False


# =========================================================
# OCR
# =========================================================

def extract_ocr_results_with_boxes(img_bytes, ocr_model, debug_name=None, small_skip_threshold=40):
    results = []

    try:
        img_np = image_bytes_to_cv2(img_bytes)
        if img_np is None:
            print(f"[DEBUG] 图片解码失败: {debug_name}")
            return results, None

        h, w = img_np.shape[:2]
        print(f"[DEBUG] {debug_name} 原始尺寸: {w}x{h}")

        if w < small_skip_threshold or h < small_skip_threshold:
            print(f"[DEBUG] {debug_name} 图片过小，跳过OCR")
            return results, img_np

        if w < 500 or h < 500:
            scale = 2
            img_np = cv2.resize(img_np, None, fx=scale, fy=scale, interpolation=cv2.INTER_CUBIC)
            print(f"[DEBUG] {debug_name} 放大后尺寸: {img_np.shape[1]}x{img_np.shape[0]}")

        ocr_raw = ocr_model.predict(img_np)
        if ocr_raw is None:
            return results, img_np

        for res in ocr_raw:
            if isinstance(res, list):
                for line in res:
                    try:
                        if isinstance(line, list) and len(line) >= 2:
                            box = normalize_box(line[0])
                            rec_info = line[1]
                            text = None
                            if isinstance(rec_info, (list, tuple)) and len(rec_info) >= 1:
                                text = str(rec_info[0]).strip()
                            if text and box:
                                results.append({"text": text, "box": box})
                    except Exception as e:
                        print(f"[DEBUG] 旧版OCR结构解析失败: {e}")

            elif isinstance(res, dict):
                rec_texts = res.get("rec_texts", None)
                rec_boxes = res.get("rec_boxes", None)
                dt_polys = res.get("dt_polys", None)

                if rec_texts is not None and rec_boxes is not None and len(rec_texts) == len(rec_boxes):
                    for text, box in zip(rec_texts, rec_boxes):
                        try:
                            text = str(text).strip()
                            box = normalize_box(box)
                            if text and box:
                                results.append({"text": text, "box": box})
                        except Exception as e:
                            print(f"[DEBUG] rec_texts + rec_boxes 解析失败: {e}")

                elif rec_texts is not None and dt_polys is not None and len(rec_texts) == len(dt_polys):
                    for text, box in zip(rec_texts, dt_polys):
                        try:
                            text = str(text).strip()
                            box = normalize_box(box)
                            if text and box:
                                results.append({"text": text, "box": box})
                        except Exception as e:
                            print(f"[DEBUG] rec_texts + dt_polys 解析失败: {e}")

        print(f"[DEBUG] {debug_name} OCR结果数量: {len(results)}")

    except Exception as e:
        print(f"[DEBUG] OCR提取文本和坐标失败 {debug_name}: {e}")
        return [], None

    return results, img_np


def draw_sensitive_areas_on_image_to_bytes(img_np, matched_items):
    if img_np is None or matched_items is None or len(matched_items) == 0:
        return None

    img_rgb = cv2.cvtColor(img_np, cv2.COLOR_BGR2RGB)

    fig, ax = plt.subplots(1, figsize=(8, 6))
    ax.imshow(img_rgb)

    for item in matched_items:
        box = item.get("box", None)
        text = item.get("text", "")

        if box is None or len(box) < 4:
            continue

        try:
            xs = [float(p[0]) for p in box]
            ys = [float(p[1]) for p in box]
            min_x, max_x = min(xs), max(xs)
            min_y, max_y = min(ys), max(ys)

            rect = patches.Rectangle(
                (min_x, min_y),
                max_x - min_x,
                max_y - min_y,
                linewidth=2,
                edgecolor='yellow',
                facecolor='yellow',
                alpha=0.35
            )
            ax.add_patch(rect)

            ax.text(
                min_x,
                min_y - 5 if min_y > 10 else min_y + 15,
                text,
                color='black',
                fontsize=8,
                bbox=dict(facecolor='yellow', alpha=0.7, edgecolor='yellow')
            )
        except Exception as e:
            print(f"[DEBUG] 绘制高亮框失败: {e}")

    plt.axis('off')
    plt.tight_layout()

    bio = io.BytesIO()
    plt.savefig(bio, format='png', bbox_inches='tight', pad_inches=0.05, dpi=150)
    plt.close(fig)
    bio.seek(0)
    return bio.read()


def process_image_and_generate_marked_bytes(img_bytes, ocr_model, simple_keywords, regex_patterns, debug_name):
    if ocr_model is None:
        ocr_model = get_ocr_model()

    matched_keywords = []
    matched_items = []

    ocr_results, img_np = extract_ocr_results_with_boxes(img_bytes, ocr_model, debug_name=debug_name)

    for item in ocr_results:
        text = item["text"]
        matches = find_matches_in_content(text, simple_keywords, regex_patterns)
        if matches:
            matched_items.append(item)
            for _, keyword, _, _ in matches:
                matched_keywords.append(keyword)

    matched_keywords = list(dict.fromkeys(matched_keywords))

    marked_img_bytes = None
    if len(matched_items) > 0:
        marked_img_bytes = draw_sensitive_areas_on_image_to_bytes(img_np, matched_items)

    return matched_keywords, marked_img_bytes


# =========================================================
# PDFLogger
# =========================================================

class PDFLogger:
    def __init__(self, filename="sensitive_info_report.pdf"):
        try:
            pdfmetrics.registerFont(TTFont('Deng', 'C:/Windows/Fonts/Deng.ttf'))
        except:
            print("警告：未找到'等线'字体，PDF中的中文可能无法显示。")

        self.filename = filename
        self.c = canvas.Canvas(filename, pagesize=letter)
        self.width, self.height = letter
        self.y_position = self.height - 50
        self.c.setFont('Deng', 18)
        self.c.drawString(50, self.height - 35, f"敏感信息扫描报告 - {datetime.now().strftime('%Y-%m-%d %H:%M')}")
        self.c.setFont('Deng', 10)

    def _check_new_page(self):
        if self.y_position < 50:
            self.c.showPage()
            self.c.setFont('Deng', 10)
            self.y_position = self.height - 50

    def add_entry(self, file_path, locations):
        self._check_new_page()
        self.c.setFont('Deng', 12)
        absolute_path = os.path.abspath(file_path)
        self.c.drawString(50, self.y_position, f"文件: {absolute_path}")
        self.y_position -= 20

        self.c.setFont('Deng', 9)
        for loc in locations:
            self._check_new_page()
            self.c.drawString(65, self.y_position, f"- {loc}")
            self.y_position -= 15
        self.y_position -= 10

    def save(self):
        self.c.save()
        print(f"\n报告已生成: {self.filename}")


# =========================================================
# DOCX
# =========================================================

def highlight_sensitive_words_in_doc(file_path, simple_keywords, regex_patterns, ocr_model):
    locations = []
    _, ext = os.path.splitext(file_path)
    found_in_file = False
    word_app = None

    try:
        if ext.lower() == '.docx':
            document = Document(file_path)

            for i, para in enumerate(document.paragraphs):
                matches = find_matches_in_content(para.text, simple_keywords, regex_patterns)
                if matches:
                    found_in_file = True
                    for _, keyword, _, _ in matches:
                        locations.append(f"第 {i+1} 段文本 (内容: '{keyword}')")
                    for run in para.runs:
                        run.font.highlight_color = WD_COLOR_INDEX.YELLOW

            image_index = 0
            for rel in document.part.rels.values():
                if "image" in rel.target_ref:
                    image_index += 1
                    img_bytes = rel.target_part.blob

                    matched_keywords, marked_img_bytes = process_image_and_generate_marked_bytes(
                        img_bytes=img_bytes,
                        ocr_model=ocr_model,
                        simple_keywords=simple_keywords,
                        regex_patterns=regex_patterns,
                        debug_name=f"docx_image_{image_index}"
                    )

                    if matched_keywords and marked_img_bytes:
                        try:
                            rel.target_part._blob = marked_img_bytes
                            found_in_file = True
                            for keyword in matched_keywords:
                                locations.append(f"图片 {image_index} 中OCR发现: '{keyword}'，已在原DOCX中标记")
                        except Exception as e:
                            locations.append(f"图片 {image_index} 命中敏感信息，但写回DOCX失败: {e}")

            if found_in_file:
                document.save(file_path)

        elif ext.lower() == '.doc':
            word_app = win32com.client.Dispatch("Word.Application")
            word_app.Visible = False
            doc = word_app.Documents.Open(os.path.abspath(file_path))

            find_obj = doc.Content.Find
            find_obj.ClearFormatting()
            find_obj.Replacement.ClearFormatting()
            find_obj.Replacement.Highlight = True

            for word in simple_keywords:
                if find_obj.Execute(FindText=word, MatchCase=False, MatchWholeWord=False, Replace=2):
                    found_in_file = True
                    locations.append(f"文档内发现关键词 '{word}' (高亮显示)")

            if found_in_file:
                doc.Save()
            doc.Close(False)

    except Exception as e:
        locations.append(f"处理 Word 文件失败: {e}")
    finally:
        if word_app:
            word_app.Quit()

    return list(dict.fromkeys(locations))


# =========================================================
# XLSX
# =========================================================

def highlight_sensitive_words_in_excel(file_path, simple_keywords, regex_patterns, ocr_model):
    locations = []
    _, ext = os.path.splitext(file_path)
    found_in_file = False
    excel_app = None

    try:
        if ext.lower() == '.xlsx':
            print(f"[DEBUG] 开始加载 Excel: {file_path}")
            workbook = openpyxl.load_workbook(file_path)
            yellow_fill = PatternFill(start_color='FFFFFF00', end_color='FFFFFF00', fill_type='solid')

            print("[DEBUG] 开始扫描单元格文本")
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value is not None:
                            matches = find_matches_in_content(str(cell.value), simple_keywords, regex_patterns)
                            if matches:
                                found_in_file = True
                                cell.fill = yellow_fill
                                for _, keyword, _, _ in matches:
                                    locations.append(f"Sheet[{sheet_name}]单元格[{cell.coordinate}] (内容: '{keyword}')")

            if found_in_file:
                workbook.save(file_path)
                print(f"[DEBUG] Excel单元格高亮已保存到原文件: {file_path}")

            replace_map = {}
            with zipfile.ZipFile(file_path, 'r') as z:
                media_names = [name for name in z.namelist() if name.startswith("xl/media/")]
                print(f"[DEBUG] 从xlsx中提取到图片数量: {len(media_names)}")

                for idx, media_name in enumerate(media_names, start=1):
                    try:
                        img_bytes = z.read(media_name)
                        matched_keywords, marked_img_bytes = process_image_and_generate_marked_bytes(
                            img_bytes=img_bytes,
                            ocr_model=ocr_model,
                            simple_keywords=simple_keywords,
                            regex_patterns=regex_patterns,
                            debug_name=f"xlsx_image_{idx}_{os.path.basename(media_name)}"
                        )

                        if matched_keywords and marked_img_bytes:
                            replace_map[media_name] = marked_img_bytes
                            found_in_file = True
                            for keyword in matched_keywords:
                                locations.append(f"Excel图片[{os.path.basename(media_name)}] 中OCR发现: '{keyword}'，已在原XLSX中标记")

                    except Exception as img_e:
                        locations.append(f"Excel图片[{media_name}] OCR失败: {img_e}")

            if len(replace_map) > 0:
                ok = replace_bytes_in_zip(file_path, replace_map)
                if ok:
                    print(f"[DEBUG] Excel图片已回写原文件: {file_path}")
                else:
                    locations.append("Excel图片命中敏感信息，但回写原文件失败")

        elif ext.lower() == '.xls':
            excel_app = win32com.client.Dispatch("Excel.Application")
            excel_app.Visible = False
            workbook = excel_app.Workbooks.Open(os.path.abspath(file_path))

            for sheet in workbook.Worksheets:
                for cell in sheet.UsedRange.Cells:
                    if cell.Value:
                        matches = find_matches_in_content(str(cell.Value), simple_keywords, regex_patterns)
                        if matches:
                            found_in_file = True
                            cell.Interior.ColorIndex = 6
                            for _, keyword, _, _ in matches:
                                locations.append(
                                    f"Sheet[{sheet.Name}] 单元格[{cell.Address.replace('$', '')}] (内容: '{keyword}')"
                                )

            if found_in_file:
                workbook.Save()
            workbook.Close(False)

    except Exception as e:
        locations.append(f"处理 Excel 文件失败: {e}")
    finally:
        if excel_app:
            excel_app.Quit()

    return list(dict.fromkeys(locations))


# =========================================================
# PPTX
# =========================================================

def highlight_sensitive_words_in_ppt(file_path, simple_keywords, regex_patterns, ocr_model):
    locations = []
    _, ext = os.path.splitext(file_path)
    found_in_file = False
    ppt_app = None
    temp_files = []

    def replace_picture_with_new_image(slide, shape, new_img_bytes):
        """
        稳定方案：删除原图，在同位置插入新图
        """
        try:
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height

            tmp_img = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
            tmp_img.write(new_img_bytes)
            tmp_img.close()
            temp_files.append(tmp_img.name)

            sp = shape._element
            sp.getparent().remove(sp)

            slide.shapes.add_picture(tmp_img.name, left, top, width=width, height=height)
            return True
        except Exception as e:
            print(f"[DEBUG] PPT图片替换失败: {e}")
            return False

    try:
        if ext.lower() == '.pptx':
            prs = Presentation(file_path)

            for slide_index, slide in enumerate(prs.slides):
                image_counter = 0

                for shape in list(slide.shapes):
                    try:
                        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                            # 组合对象这里先跳过，避免复杂结构导致报错
                            continue

                        if hasattr(shape, "has_text_frame") and shape.has_text_frame and shape.text_frame.text:
                            matches = find_matches_in_content(shape.text_frame.text, simple_keywords, regex_patterns)
                            if matches:
                                found_in_file = True
                                for _, keyword, _, _ in matches:
                                    locations.append(f"第 {slide_index+1} 页文本框 (内容: '{keyword}')")
                                for para in shape.text_frame.paragraphs:
                                    for run in para.runs:
                                        run.font.color.rgb = PptxRGBColor(255, 0, 0)

                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            image_counter += 1
                            try:
                                img_bytes = shape.image.blob
                            except Exception as e:
                                locations.append(f"第 {slide_index+1} 页图片 {image_counter} 读取失败: {e}")
                                continue

                            matched_keywords, marked_img_bytes = process_image_and_generate_marked_bytes(
                                img_bytes=img_bytes,
                                ocr_model=ocr_model,
                                simple_keywords=simple_keywords,
                                regex_patterns=regex_patterns,
                                debug_name=f"ppt_slide_{slide_index+1}_image_{image_counter}"
                            )

                            if matched_keywords and marked_img_bytes:
                                ok = replace_picture_with_new_image(slide, shape, marked_img_bytes)
                                if ok:
                                    found_in_file = True
                                    for keyword in matched_keywords:
                                        locations.append(f"第 {slide_index+1} 页图片 {image_counter} 中OCR发现: '{keyword}'，已在原PPTX中标记")
                                else:
                                    for keyword in matched_keywords:
                                        locations.append(f"第 {slide_index+1} 页图片 {image_counter} 命中敏感信息，但写回PPT失败")

                    except Exception as shape_e:
                        locations.append(f"第 {slide_index+1} 页某对象处理失败: {shape_e}")

            if found_in_file:
                prs.save(file_path)

        elif ext.lower() == '.ppt':
            ppt_app = win32com.client.Dispatch("PowerPoint.Application")
            ppt_app.Visible = False
            presentation = ppt_app.Presentations.Open(os.path.abspath(file_path), WithWindow=False)

            for i, slide in enumerate(presentation.Slides):
                for shape in slide.Shapes:
                    if shape.HasTextFrame and shape.TextFrame.HasText:
                        text_range = shape.TextFrame.TextRange
                        text_content = text_range.Text
                        matches = find_matches_in_content(str(text_content), simple_keywords, regex_patterns)
                        if matches:
                            found_in_file = True
                            for _, keyword, _, _ in matches:
                                locations.append(f"第 {i+1} 页幻灯片文本框 (内容: '{keyword}')")
                            text_range.Font.Color.RGB = 255

            if found_in_file:
                presentation.Save()
            presentation.Close()

    except Exception as e:
        locations.append(f"处理 PowerPoint 文件失败: {e}")
    finally:
        if ppt_app:
            ppt_app.Quit()

        for f in temp_files:
            try:
                os.remove(f)
            except:
                pass

    return list(dict.fromkeys(locations))


# =========================================================
# PDF
# =========================================================

def highlight_sensitive_words_in_pdf(file_path, simple_keywords, regex_patterns, ocr_model):
    locations = []
    doc = None
    pdf_changed = False

    try:
        doc = fitz.open(file_path)

        for page_index in range(len(doc)):
            page = doc[page_index]
            img_list = page.get_images(full=True)

            if not img_list:
                continue

            # 每页最多处理前3张图，进一步减轻卡顿
            img_list = img_list[:3]
            print(f"[DEBUG] PDF第 {page_index+1} 页图片数量(处理前3张): {len(img_list)}")

            page_rects_done = set()

            for img_idx, img_info in enumerate(img_list, start=1):
                try:
                    xref = img_info[0]
                    base_image = doc.extract_image(xref)
                    if not base_image:
                        continue

                    img_bytes = base_image["image"]
                    img_np = image_bytes_to_cv2(img_bytes)
                    if img_np is None:
                        continue

                    h, w = img_np.shape[:2]
                    if w < 50 or h < 50:
                        print(f"[DEBUG] PDF第 {page_index+1} 页图片 {img_idx} 太小，跳过")
                        continue

                    # PDF这里不生成plt图片，只做OCR命中判断，避免额外开销
                    matched_keywords, _ = process_image_and_generate_marked_bytes(
                        img_bytes=img_bytes,
                        ocr_model=ocr_model,
                        simple_keywords=simple_keywords,
                        regex_patterns=regex_patterns,
                        debug_name=f"pdf_page_{page_index+1}_img_{img_idx}"
                    )

                    if not matched_keywords:
                        continue

                    rects = page.get_image_rects(xref)
                    if not rects:
                        continue

                    for rect in rects:
                        rect_key = (round(rect.x0, 2), round(rect.y0, 2), round(rect.x1, 2), round(rect.y1, 2))
                        if rect_key in page_rects_done:
                            continue
                        page_rects_done.add(rect_key)

                        annot = page.add_rect_annot(rect)
                        annot.set_colors(stroke=(1, 1, 0), fill=(1, 1, 0))
                        annot.set_border(width=2)
                        annot.set_opacity(0.20)
                        annot.update()
                        pdf_changed = True

                    for keyword in matched_keywords:
                        locations.append(f"PDF第 {page_index+1} 页图片 {img_idx} 中OCR发现: '{keyword}'，已在原PDF中标记")

                except Exception as img_e:
                    locations.append(f"PDF第 {page_index+1} 页图片 {img_idx} 处理失败: {img_e}")

        if pdf_changed:
            temp_pdf = file_path + ".tmp.pdf"
            doc.save(temp_pdf, garbage=4, deflate=True)
            doc.close()
            doc = None
            os.remove(file_path)
            shutil.move(temp_pdf, file_path)
            print(f"[DEBUG] PDF标记已写回原文件: {file_path}")

    except Exception as e:
        locations.append(f"处理 PDF 文件失败: {e}")
    finally:
        if doc:
            doc.close()

    return list(dict.fromkeys(locations))


# =========================================================
# TXT / CSV
# =========================================================

def check_sensitive_words_in_text(file_path, simple_keywords, regex_patterns):
    locations = []
    encodings = ['utf-8-sig', 'gbk', 'utf-8']

    for encoding in encodings:
        try:
            with open(file_path, 'r', encoding=encoding) as f:
                for i, line in enumerate(f):
                    line_for_regex = line.replace(',', ' ').replace(';', ' ')
                    matches = find_matches_in_content(line, simple_keywords, [])
                    matches.extend(find_matches_in_content(line_for_regex, [], regex_patterns))

                    for _, keyword, start, end in matches:
                        locations.append(f"第 {i+1} 行, 第 {start+1}-{end} 列 (内容: '{keyword.strip()}')")
            return locations

        except UnicodeDecodeError:
            continue
        except Exception as e:
            locations.append(f"处理文本文件失败: {e}")
            return locations

    return locations


# =========================================================
# 总调度
# =========================================================
def build_process_result(file_path, locations):
    return {
        "file_path": os.path.abspath(file_path),
        "matches": [{"value": loc} for loc in locations]
    }


def process_file(file_path, simple_keywords, regex_patterns, ocr_model):
    if not os.path.exists(file_path):
        print(f"错误：文件 '{file_path}' 不存在。")
        return build_process_result(file_path, [])

    _, ext = os.path.splitext(file_path)
    ext = ext.lower()
    locations = []

    print(f"--- 正在处理: {os.path.basename(file_path)} ---")

    try:
        if ext in ['.doc', '.docx']:
            locations = highlight_sensitive_words_in_doc(file_path, simple_keywords, regex_patterns, ocr_model)
        elif ext in ['.xls', '.xlsx']:
            locations = highlight_sensitive_words_in_excel(file_path, simple_keywords, regex_patterns, ocr_model)
        elif ext in ['.ppt', '.pptx']:
            locations = highlight_sensitive_words_in_ppt(file_path, simple_keywords, regex_patterns, ocr_model)
        elif ext in ['.txt', '.csv']:
            locations = check_sensitive_words_in_text(file_path, simple_keywords, regex_patterns)
        elif ext == '.pdf':
            locations = highlight_sensitive_words_in_pdf(file_path, simple_keywords, regex_patterns, ocr_model)
        else:
            print(f"不支持的文件类型: {ext}")
            return build_process_result(file_path, [])
    except Exception as e:
        print(f"处理文件失败: {file_path}, 错误: {e}")
        locations = [f"处理失败: {e}"]

    locations = list(dict.fromkeys(locations))

    if locations:
        print(f"文件 '{os.path.basename(file_path)}' 中发现 {len(locations)} 处敏感内容:")
        for loc in locations:
            print(f"  - {loc}")
    else:
        print(f"文件 '{os.path.basename(file_path)}' 中未发现敏感字符。")

    print("-" * (len(os.path.basename(file_path)) + 20) + "\n")
    return build_process_result(file_path, locations)


def manage_keywords(action: str, keyword: str = None):
    global SIMPLE_KEYWORDS

    if action == 'add':
        if keyword and keyword.strip():
            cleaned = keyword.strip()
            if cleaned not in SIMPLE_KEYWORDS:
                SIMPLE_KEYWORDS.append(cleaned)
                print(f"已添加敏感词: '{cleaned}'")
            else:
                print(f"敏感词 '{cleaned}' 已存在，无需重复添加")
        else:
            print("添加时 keyword 不能为空")

    elif action == 'remove':
        if keyword and keyword.strip():
            cleaned = keyword.strip()
            if cleaned in SIMPLE_KEYWORDS:
                SIMPLE_KEYWORDS.remove(cleaned)
                print(f"已删除敏感词: '{cleaned}'")
            else:
                print(f"敏感词 '{cleaned}' 不存在")
        else:
            print("删除时 keyword 不能为空")

    elif action == 'view':
        print("当前敏感词列表（共 {} 个）：".format(len(SIMPLE_KEYWORDS)))
        for i, kw in enumerate(SIMPLE_KEYWORDS, 1):
            print(f"  {i:2d}. {kw}")
    else:
        print("不支持的操作，请使用 'add'、'remove' 或 'view'")


def log(msg):
    try:
        base_dir = get_base_dir()
        log_path = os.path.join(base_dir, "explorer.log")
        with open(log_path, "a", encoding="utf-8") as f:
            f.write(msg + "\n")
    except Exception:
        pass


def get_base_dir():
    """获取程序运行根目录（兼容 PyInstaller）。"""
    if getattr(sys, "frozen", False):
        return os.path.dirname(sys.executable)
    return os.path.dirname(os.path.abspath(__file__))


def get_file_owner(filename):
    """获取文件所有者（需要 pywin32）。"""
    sd = win32security.GetFileSecurity(filename, win32security.OWNER_SECURITY_INFORMATION)
    owner_sid = sd.GetSecurityDescriptorOwner()
    name, _, _ = win32security.LookupAccountSid(None, owner_sid)
    return name


def get_windows_drives():
    """获取 Windows 系统所有可用的盘符。"""
    drives = []
    for letter in string.ascii_uppercase:
        drive = f"{letter}:\\"
        if os.path.exists(drive):
            drives.append(drive)
    return drives


def get_resource_base_dir():
    """兼容 PyInstaller 资源目录。"""
    if getattr(sys, "frozen", False):
        return getattr(sys, "_MEIPASS", os.path.dirname(sys.executable))
    return os.path.dirname(os.path.abspath(__file__))


def get_resource_path(relative_path):
    return os.path.join(get_resource_base_dir(), relative_path)


def get_model_home():
    return get_resource_path(os.path.join("models", "paddlex_models"))


def ensure_paddle_model_home():
    model_home = get_model_home()
    os.environ["PADDLE_PDX_MODEL_HOME"] = model_home
    return model_home


def get_ocr_model():
    global OCR_MODEL
    global OCR_INIT_ERROR

    if OCR_MODEL is not None:
        return OCR_MODEL
    if OCR_INIT_ERROR is not None:
        raise RuntimeError(f"OCR model init failed: {OCR_INIT_ERROR}")

    model_home = ensure_paddle_model_home()
    if not os.path.exists(model_home):
        OCR_INIT_ERROR = f"model directory not found: {model_home}"
        raise FileNotFoundError(OCR_INIT_ERROR)

    print(f"初始化 OCR，模型目录: {model_home}")
    log(f"初始化 OCR，模型目录: {model_home}")

    try:
        OCR_MODEL = PaddleOCR(
            use_textline_orientation=True,
            lang="ch",
            enable_mkldnn=False
        )
        return OCR_MODEL
    except Exception as e:
        OCR_INIT_ERROR = str(e)
        raise


def load_rules_from_file(file_path):
    """
    从文本文件加载规则，每行一条。
    自动忽略空行和以 # 开头的注释行。
    """
    if not file_path:
        return []

    if not os.path.exists(file_path):
        raise FileNotFoundError(f"规则文件不存在: {file_path}")

    encodings = ["utf-8-sig", "utf-8", "gbk"]
    last_error = None

    for encoding in encodings:
        try:
            rules = []
            with open(file_path, "r", encoding=encoding) as f:
                for line in f:
                    line = line.strip()
                    if not line or line.startswith("#"):
                        continue
                    rules.append(line)
            return rules
        except Exception as e:
            last_error = e

    raise RuntimeError(f"读取规则文件失败: {file_path}, 错误: {last_error}")


def merge_unique_rules(default_rules, extra_rules):
    merged = list(default_rules)
    for rule in extra_rules:
        if rule not in merged:
            merged.append(rule)
    return merged


def parse_args():
    parser = argparse.ArgumentParser(description="敏感信息扫描工具")
    parser.add_argument(
        "scan_path",
        nargs="?",
        default=None,
        help="待扫描目录或文件路径，默认扫描程序目录下的 file 文件夹"
    )
    parser.add_argument(
        "--output",
        default=None,
        help="结果 JSON 输出路径，默认输出到程序根目录下 result.json"
    )
    parser.add_argument(
        "--keywords-file",
        default=None,
        help="关键词文件路径，每行一个关键词"
    )
    parser.add_argument(
        "--regex-file",
        default=None,
        help="正则文件路径，每行一个正则表达式"
    )
    return parser.parse_args()


if __name__ == '__main__':
    try:
        args = parse_args()

        SIMPLE_KEYWORDS = [
            "监狱", "救我", "万元", "亿元", "密码", "password", "API_KEY", "token", "help me",
            "-----BEGIN PRIVATE KEY-----", "jdbc:", "mysql:", "绝密", "机密", "秘密",
            "内部", "限制分发", "工资表", "财务", "报表", "账户", "名单", "投标", "报价",
            "地址", "address"
        ]

        REGEX_PATTERNS = [
            r'\b[1-9]\d{5}(?:18|19|20)\d{2}(?:0[1-9]|1[0-2])(?:0[1-9]|[12]\d|3[01])\d{3}[\dX]\b',
            r'\b[1-9]\d{7}(?:0[1-9]|1[0-2])(?:0[1-9]|[12]\d|3[01])\d{3}\b',
            r'\b[GESPDgespd]\d{8}\b',
            r'[\u4e00-\u9fa5]字第\d{5,8}号',
            r'\b1[3-9]\d{9}\b',
            r'\b\d{3,4}-\d{7,8}\b',
            r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b',
            r'\b\d{16,19}\b',
            r'\b[A-Z0-9]{18}\b',
            r'\b[A-Z0-9]{15,20}\b',
            r'\b\d{8}\b|\b\d{10}\b',
            r'\b(?:10\.\d{1,3}\.\d{1,3}\.\d{1,3}|172\.(?:1[6-9]|2\d|3[01])\.\d{1,3}\.\d{1,3}|192\.168\.\d{1,3}\.\d{1,3})\b',
            r'\b(?:[0-9A-Fa-f]{2}[:-]){5}[0-9A-Fa-f]{2}\b',
            r'\b[京津沪渝冀豫云辽黑湘皖鲁新苏浙赣鄂桂甘晋蒙陕吉闽贵粤青藏川宁琼使领][A-Z](?:[A-HJ-NP-Z0-9]{5}|[A-HJ-NP-Z0-9]{4}[DF])\b',
            r'\b[A-HJ-NPR-Z0-9]{17}\b',
        ]

        # 从外部文件追加关键词 / 正则
        extra_keywords = load_rules_from_file(args.keywords_file) if args.keywords_file else []
        extra_regex = load_rules_from_file(args.regex_file) if args.regex_file else []

        SIMPLE_KEYWORDS = merge_unique_rules(SIMPLE_KEYWORDS, extra_keywords)
        REGEX_PATTERNS = merge_unique_rules(REGEX_PATTERNS, extra_regex)

        print(f"关键词数量: {len(SIMPLE_KEYWORDS)}")
        print(f"正则数量: {len(REGEX_PATTERNS)}")
        log(f"关键词数量: {len(SIMPLE_KEYWORDS)}")
        log(f"正则数量: {len(REGEX_PATTERNS)}")

        print("初始化 OCR...")
        log("初始化 OCR...")
        os.environ["PADDLE_PDX_DISABLE_MODEL_SOURCE_CHECK"] = "True"

        ensure_paddle_model_home()

        BASE_DIR = get_base_dir()
        SCAN_PATH = args.scan_path if args.scan_path else os.path.join(BASE_DIR, "file")
        output_path = args.output if args.output else os.path.join(BASE_DIR, "result.json")

        print(f"扫描路径: {SCAN_PATH}")
        print(f"输出路径: {output_path}")
        log(f"扫描路径: {SCAN_PATH}")
        log(f"输出路径: {output_path}")

        if not os.path.exists(SCAN_PATH):
            raise FileNotFoundError(f"扫描路径不存在: {SCAN_PATH}")

        all_files = []

        if os.path.isfile(SCAN_PATH):
            if SCAN_PATH.lower().endswith(TARGET_EXTENSIONS) and not os.path.basename(SCAN_PATH).startswith(SKIP_FILE_PREFIXES):
                all_files.append(os.path.abspath(SCAN_PATH))
        else:
            for root, _, files in os.walk(SCAN_PATH):
                for file in files:
                    if file.startswith(SKIP_FILE_PREFIXES):
                        continue
                    if file.lower().endswith(TARGET_EXTENSIONS):
                        all_files.append(os.path.join(root, file))

        print(f"共找到 {len(all_files)} 个文件")
        log(f"共找到 {len(all_files)} 个文件")

        final_results = []

        for file_path in tqdm(all_files, desc="识别中"):
            try:
                result = process_file(
                    file_path,
                    SIMPLE_KEYWORDS,
                    REGEX_PATTERNS,
                    None
                )

                if result.get("matches"):
                    final_results.append(result)

            except Exception as e:
                err = f"处理失败: {file_path}, 错误: {e}"
                print(err)
                log(err)
                log(traceback.format_exc())

        os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(final_results, f, ensure_ascii=False, indent=2)

        print(f"\n识别完成，结果已保存到: {output_path}")
        log(f"识别完成，结果已保存到: {output_path}")

    except Exception as e:
        err = f"程序启动失败: {e}"
        print(err)
        log(err)
        log(traceback.format_exc())
        raise
