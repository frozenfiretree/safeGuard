import logging
from pathlib import Path
from zipfile import BadZipFile

from docx import Document
import pdfplumber
import openpyxl
from pptx import Presentation

from .chunker import TextBlock, ContentLocation


logger = logging.getLogger(__name__)


def txt_extractor(path):
    text = Path(path).read_text(encoding="utf-8", errors="ignore")
    return [text]


def _safe_text_fallback(path, reason: str):
    logger.warning("extractor fallback to text mode: path=%s, reason=%s", path, reason)
    texts = txt_extractor(path)
    return [t for t in texts if t and t.strip()]


def docx_extractor(path):
    try:
        doc = Document(path)
        return [p.text for p in doc.paragraphs if p.text.strip()]
    except (ValueError, BadZipFile, KeyError) as e:
        return _safe_text_fallback(path, f"docx_parse_failed:{e}")


def pdf_extractor(path):
    texts = []
    try:
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    texts.append(t)
    except Exception as e:
        return _safe_text_fallback(path, f"pdf_parse_failed:{e}")
    return texts


def xlsx_extractor(path):
    try:
        wb = openpyxl.load_workbook(path, data_only=True)

        texts = []
        for sheet in wb:
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join(str(c) for c in row if c)
                if row_text:
                    texts.append(row_text)

        return texts
    except (ValueError, BadZipFile, KeyError) as e:
        return _safe_text_fallback(path, f"xlsx_parse_failed:{e}")


def pptx_extractor(path):
    try:
        prs = Presentation(path)

        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text.strip():
                        texts.append(shape.text)

        return texts
    except (ValueError, BadZipFile, KeyError) as e:
        return _safe_text_fallback(path, f"pptx_parse_failed:{e}")


EXTRACTORS = {
    ".txt": txt_extractor,
    ".docx": docx_extractor,
    ".pdf": pdf_extractor,
    ".xlsx": xlsx_extractor,
    ".pptx": pptx_extractor,
}


def file_to_textblocks(file_path):
    file_path = Path(file_path)

    extractor = EXTRACTORS.get(file_path.suffix.lower())

    if extractor is None:
        return []

    texts = extractor(file_path)

    blocks = []

    for i, text in enumerate(texts):

        blocks.append(
            TextBlock(
                content=text,
                source_type="paragraph",
                location=ContentLocation(
                    block_index=i,
                    char_offset_start=0,
                    char_offset_end=len(text),
                ),
            )
        )

    return blocks
